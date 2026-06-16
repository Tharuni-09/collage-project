from flask import Flask, send_file, render_template, request, redirect, session, current_app, jsonify
from werkzeug.utils import secure_filename
# import google.generativeai as genai
import os
import logging
import time
import sqlite3
import hashlib
import re
import io
import json
import requests
from flask import url_for
from datetime import datetime
from reportlab.lib.pagesizes import A4
from reportlab.platypus import (
    SimpleDocTemplate,
    BaseDocTemplate,
    PageTemplate,
    Frame,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    HRFlowable,
    FrameBreak,
    Image,
    KeepTogether,
    PageBreak
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY, TA_RIGHT
from reportlab.pdfbase import pdfform
from reportlab.lib.units import mm
from pptx import Presentation
from reportlab.platypus import *
import traceback
# import google.generativeai as genai_v1

load_dotenv()

# genai.configure(api_key=GEMINI_API_KEY)
# model = genai.GenerativeModel("gemini-1.5-flash")

# model = genai_v1.GenerativeModel(
#     "gemini-1.5-flash"
# )

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, "database", "ml_dept.db")

GENERATED_FOLDER = os.path.join(BASE_DIR, 'static', 'uploads', 'generated')
os.makedirs(GENERATED_FOLDER, exist_ok=True)

NOTES_UPLOAD_FOLDER = os.path.join(BASE_DIR, 'static', 'uploads', 'notes')
os.makedirs(NOTES_UPLOAD_FOLDER, exist_ok=True)
ALLOWED_NOTES_EXTENSIONS = {'pdf', 'ppt', 'pptx'}

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# Configure logging to help diagnose connection issues
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Get API Key from environment variable for security
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")

try:
    client = genai.Client(
        api_key=GEMINI_API_KEY
    )
except Exception as e:
    logger.error(f"Failed to initialize Gemini AI client: {e}")
    client = None


# =====================================================================
# HELPER FUNCTIONS
# =====================================================================

def extract_text_from_file(filepath):
    """Helper to extract text from TXT or PDF files."""
    ext = filepath.rsplit('.', 1)[1].lower()
    if ext == 'txt':
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"TXT extraction failed: {e}")
            return ""
    elif ext == 'pdf':
        try:
            from pypdf import PdfReader
            reader = PdfReader(filepath)
            text = ""
            for page in reader.pages:
                content = page.extract_text()
                if content:
                    text += content + "\n"
            return text
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return ""
    return ""

def create_ppt(source_filepath, output_filename, slides_count):
    """Implementation of PPT generation using Gemini AI."""
    text = extract_text_from_file(source_filepath)
    content_data = []

    if client and text:
        prompt = (
            f"Summarize the following text into {slides_count} presentation slides. "
            "For each slide, provide a Title, 3 bullet points, and one specific keyword for a relevant stock image. "
            "Format your response exactly like this for each slide:\n"
            "Slide: [Title]\n"
            "* [Bullet 1]\n"
            "* [Bullet 2]\n"
            "* [Bullet 3]\n\n"
            "Keyword: [Single Word]\n\n"
            f"Text: {text[:8000]}"
        )
        try:
            response = response = model.generate_content(prompt)
            raw_text = response.text
            slides_raw = raw_text.split("Slide:")
            for s in slides_raw:
                if not s.strip(): continue
                lines = s.strip().split('\n')
                title = lines[0].strip()
                bullets = [l.strip('* ').strip() for l in lines[1:] if l.strip().startswith('*')]
                keyword = "technology"
                for line in lines:
                    if line.startswith("Keyword:"):
                        keyword = line.replace("Keyword:", "").strip()
                if title:
                    content_data.append({"title": title, "bullets": bullets[:3], "keyword": keyword})
        except Exception as e:
            logger.error(f"AI PPT Generation failed: {e}")

    # Fallback to skeleton if AI fails or no content generated
    if not content_data:
        for i in range(int(slides_count)):
            content_data.append({
                "title": f"Topic {i+1} from {os.path.basename(source_filepath)}",
                "bullets": ["Key takeaway point 1", "Key takeaway point 2", "Key takeaway point 3"],
                "keyword": "education"
            })

    prs = Presentation()
    slide_layout = prs.slide_layouts[1] # Title and Content
    
    for data in content_data[:int(slides_count)]:
        slide = prs.slides.add_slide(slide_layout)
        
        # Add Image if possible
        try:
            img_url = f"https://source.unsplash.com/featured/400x300?{data['keyword']}"
            img_data = requests.get(img_url, timeout=5).content
            img_temp = io.BytesIO(img_data)
            # Place image in the top right corner
            from pptx.util import Inches
            slide.shapes.add_picture(img_temp, Inches(6.5), Inches(1), width=Inches(3))
        except Exception as e:
            logger.warning(f"Could not add image for slide '{data['title']}': {e}")

        slide.shapes.title.text = data["title"]
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        for i, bullet in enumerate(data["bullets"]):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet

    output_path = os.path.join(GENERATED_FOLDER, output_filename)
    prs.save(output_path)
    return output_path, content_data

def create_notes_pdf(source_filepath, output_filename, section_count):
    """Implementation of Notes generation using Gemini AI and ReportLab."""
    text = extract_text_from_file(source_filepath)
    content_data = []

    if client and text:
        prompt = (
            f"Generate detailed study notes in {section_count} sections for this text. "
            "For each section, provide a Heading and a Summary paragraph. "
            "Format exactly as:\nSection: [Heading]\nSummary: [Paragraph]\n\n"
            f"Text: {text[:8000]}"
        )
        try:
            response = client.models.generate_content(model="gemini-1.5-flash", contents=prompt)
            sections_raw = response.text.split("Section:")
            for s in sections_raw:
                if "Summary:" in s:
                    parts = s.split("Summary:")
                    content_data.append({"heading": parts[0].strip(), "summary": parts[1].strip()})
        except Exception as e:
            logger.error(f"AI Notes Generation failed: {e}")

    if not content_data:
        for i in range(int(section_count)):
            content_data.append({"heading": f"Section {i+1}", "summary": "Detailed analysis pending AI processing."})

    output_path = os.path.join(GENERATED_FOLDER, output_filename)
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(f"Study Notes: {os.path.basename(source_filepath)}", styles['Title']), Spacer(1, 12)]
    for data in content_data[:int(section_count)]:
        story.append(Paragraph(data["heading"], styles['Heading2']))
        story.append(Paragraph(data["summary"], styles['Normal']))
        story.append(Spacer(1, 12))
    doc.build(story)
    return output_path, content_data

def should_print_section(text, user_response):
    if not text or not text.strip():
        return False
    low = user_response.strip().lower()
    if low in ["no", "false", "0", "none", "not applicable", "na"]:
        return False
    return True


def elaborate(text, section_name):
    if not text or section_name not in ["skills", "projects", "internships"]:
        return text

    lines = [line.strip() for line in text.split("\n") if line.strip()]
    if not lines:
        return ""

    if section_name == "skills":
        bullets = []
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.lower().startswith("python"):
                bullets.append(
                    "Python: Proficient in core Python, OOP, and popular libraries such as NumPy and Pandas."
                )
            elif line.lower().startswith("flask"):
                bullets.append(
                    "Flask: Developed web applications and REST APIs using Flask and SQLite."
                )
            elif line.lower().startswith("machine learning") or "ml" in line.lower():
                bullets.append(
                    "Machine Learning: Experience with classification, regression, and neural networks using scikit‑learn or TensorFlow/PyTorch."
                )
            elif "sql" in line.lower() or "database" in line.lower():
                bullets.append(
                    "SQL & Databases: Designed and queried SQLite and MySQL databases for web applications."
                )
            elif "html" in line.lower() and "css" in line.lower():
                bullets.append(
                    "Web Technologies: Built responsive websites using HTML, CSS, and JavaScript (or frameworks)."
                )
            else:
                bullets.append(f"{line}: Comfortable applying this technology to real‑world projects.")
        return "\n".join(bullets)

    if section_name == "projects":
        bullets = []
        for i, line in enumerate(lines):
            if i == 0 and len(lines) == 1:
                bullets.append(
                    f"{line} – A well‑structured project applying core concepts of the technology stack."
                )
            elif i == 0 and len(lines) > 1:
                bullets.append(f"Project: {line}")
            else:
                bullets.append(f"• {line}")
        return "\n".join(bullets)

    if section_name == "internships":
        bullets = []
        for line in lines:
            if line.strip().lower().startswith("intern") or "internship" in line.lower():
                bullets.append(
                    "Internship role focused on practical software development in a team environment."
                )
            else:
                bullets.append(f"• {line}")
        return "\n".join(bullets)

    return text


def polish_objective(obj):
    if not obj.strip():
        return ""
    words = obj.strip().split()
    if not words:
        return ""
    obj = " ".join(words)
    if not obj.endswith("."):
        obj += "."
    return f"Motivated and detail‑oriented professional with strong {obj[0].lower() + obj[1:]}"


# Expands user input into a fuller, one‑page‑style section
def expand_section(text, section_name):
    if not text or text.lower().strip() in ["no", "none", "na", "false", "0"]:
        return ""

    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if not lines:
        return ""

    expanded = []

    if section_name == "objective":
        return (
            "Accomplished professional with a background in "
            f"{text}. Dedicated to leveraging technical skills to drive operational "
            "excellence and contribute to organizational success in a challenging role."
        )

    elif section_name == "work_experience":
        for line in lines:
            expanded.append(f"• {line}")
            expanded.append(
                "  - Achieved key deliverables by optimizing workflows and improving "
                "efficiency by 15–20%."
            )
            expanded.append(
                "  - Collaborated with cross‑functional teams to resolve complex "
                "technical challenges."
            )
        return "\n".join(expanded)

    elif section_name == "projects":
        for line in lines:
            expanded.append(f"• Project: {line}")
            expanded.append(
                "  - Developed core architecture using modern frameworks, ensuring "
                "scalability and high performance."
            )
            expanded.append(
                "  - Conducted thorough testing and debugging, resulting in a 10% "
                "improvement in system stability."
            )
        return "\n".join(expanded)

    return "\n".join([f"• {l}" for l in lines])


COMMON_SPELLING_FIXES = {
    "oblective": "objective",
    "doutd": "doubt",
    "recieve": "receive",
    "teh": "the",
    "adn": "and",
    "acheivement": "achievement",
    "deatils": "details",
    "develpment": "development",
    "manger": "manager",
    "projct": "project",
    "programe": "program",
    "speling": "spelling",
    "currect": "correct",
}


def clean_text(text):
    if not text:
        return ""
    # Clean and escape for ReportLab Paragraphs
    text = str(text).strip()
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    text = text.replace("\n", "<br/>")
    for wrong, right in COMMON_SPELLING_FIXES.items():
        text = re.sub(rf"\b{wrong}\b", right, text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text)
    return text

def allowed_notes_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_NOTES_EXTENSIONS

def split_items(text):
    if not text:
        return []
    # Split by newline, comma, or semicolon and remove whitespace
    return [i.strip() for i in re.split(r'[\n,;]+', str(text)) if i.strip()]


def enhance_objective(objective, field_of_study):
    objective = clean_text(objective)
    if not objective:
        field = field_of_study or "Machine Learning"
        return f"Results-oriented {field} professional with a focus on delivering high-impact technical solutions through advanced machine learning methodologies and scalable software architecture."
    if len(objective.split()) < 15:
        objective = objective.rstrip('.')
        objective += ", eager to contribute to innovative teams and dedicated to leveraging data-driven insights to solve complex challenges."
    return objective


def parse_projects(projects_text):
    if not projects_text or not projects_text.strip():
        return []
    return [clean_text(line) for line in projects_text.splitlines() if line.strip()]

def describe_project(project_line):
    return []

def describe_experience(exp_line):
    text = exp_line.strip()
    lower = text.lower()
    bullets = []
    
    if any(w in lower for w in ["software", "developer", "sde", "engineer", "backend", "frontend"]):
        bullets.append(f"Contributed to high-quality codebases as a {text}, focusing on scalable features and system reliability.")
        bullets.append("Collaborated with cross-functional teams to deliver robust software solutions using modern frameworks and best practices.")
        bullets.append("Optimized application performance and resolved complex technical bottlenecks in production environments.")
    elif any(w in lower for w in ["data", "analyst", "scientist", "ml", "machine learning"]):
        bullets.append(f"Leveraged statistical methodologies in my role as {text} to extract actionable insights from complex datasets.")
        bullets.append("Developed automated reporting pipelines and visualization dashboards to drive data-driven decision making.")
        bullets.append("Performed deep-dive exploratory data analysis to identify trends and improve model predictive accuracy.")
    elif any(w in lower for w in ["research", "intern", "academic"]):
        bullets.append(f"Conducted in-depth research for {text}, contributing to critical project milestones and technical documentation.")
        bullets.append("Documented technical findings and presented comprehensive reports to key stakeholders and mentors.")
        bullets.append("Assisted in the design and execution of experimental workflows to validate theoretical hypotheses.")
    else:
        bullets.append(f"Supported organizational goals as {text} through effective task management and cross-team communication.")
        bullets.append("Developed professional domain expertise while contributing to organizational efficiency and project success.")
        bullets.append("Maintained high standards of quality and reliability in all assigned deliverables.")
        
    return bullets

def describe_certification(cert_line):
    text = cert_line.strip()
    if not text: return ""
    lower = text.lower()
    if any(w in lower for w in ["aws", "google", "azure", "cloud"]):
        return f"Validated cloud architecture expertise and proficiency in designing secure, scalable distributed systems."
    if any(w in lower for w in ["python", "java", "sql", "javascript"]):
        return f"Demonstrated advanced technical proficiency and problem-solving capabilities in {text}."
    return f"Professional credential validating advanced proficiency and theoretical knowledge in {text}."

# =====================================================================
# APP SETUP
# =====================================================================

app = Flask(
    __name__,
    static_folder=os.path.join(os.path.dirname(__file__), "static")
)
app.secret_key = "your_secret_key_here"

def get_db():
    os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)  # Auto-create folder
    conn = sqlite3.connect(DB_PATH, timeout=10)
    conn.row_factory = sqlite3.Row
    return conn


# =====================================================================
# INIT DB SCHEMA
# =====================================================================

def init_db():
    """Ensures the database is initialized with the required schema."""
    # Check if the 'users' table exists to decide if we need to run schema.sql
    db = get_db()
    try:
        table_exists = db.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='users'").fetchone()
        if not table_exists:
            with app.open_resource("schema.sql") as f:
                db.executescript(f.read().decode("utf-8"))

        # Ensure is_hod column exists (Migration support)
        try: db.execute("ALTER TABLE users ADD COLUMN is_hod INTEGER DEFAULT 0")
        except: pass
        
        # Ensure branch column exists in previous_year_papers (Migration support)
        try: db.execute("ALTER TABLE previous_year_papers ADD COLUMN branch TEXT DEFAULT 'ALL'")
        except: pass
        
        # Ensure photo_limit column exists in previous_year_papers (Migration support)
        try: db.execute("ALTER TABLE previous_year_papers ADD COLUMN photo_limit INTEGER DEFAULT 1")
        except: pass
        
        # Ensure previous_year_papers table exists
        db.execute("""
            CREATE TABLE IF NOT EXISTS previous_year_papers (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                subject TEXT NOT NULL,
                branch TEXT NOT NULL,
                year INTEGER NOT NULL,
                semester INTEGER NOT NULL,
                photo_limit INTEGER DEFAULT 1 -- Serves as No. of Pages
            )
        """)
        db.execute("""
            CREATE TABLE IF NOT EXISTS paper_photos (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                paper_id INTEGER NOT NULL,
                image_path TEXT NOT NULL,
                FOREIGN KEY (paper_id) REFERENCES previous_year_papers(id)
            )
        """)
        # Ensure notes and note_files tables exist
        db.execute("""
            CREATE TABLE IF NOT EXISTS notes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                subject TEXT NOT NULL,
                semester INTEGER NOT NULL,
                department TEXT NOT NULL,
                year INTEGER NOT NULL,
                uploaded_by_uid INTEGER NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (uploaded_by_uid) REFERENCES users(uid)
            )
        """)
        db.execute("""
            CREATE TABLE IF NOT EXISTS note_files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                note_id INTEGER NOT NULL,
                filename TEXT NOT NULL,
                file_path TEXT NOT NULL,
                file_type TEXT NOT NULL,
                FOREIGN KEY (note_id) REFERENCES notes(id)
            )
        """)
        # Ensure generated_files table exists
        db.execute("""
            CREATE TABLE IF NOT EXISTS generated_files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_name TEXT NOT NULL,
                generated_file TEXT NOT NULL,
                file_type TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                slide_json TEXT -- Added slide_json column
            )
        """)
        db.commit()
    finally:
        db.close()

# =====================================================================
# INDEX AND STATIC PAGES
# =====================================================================

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/outreach")
def outreach():
    return render_template("outreach.html")


@app.route("/gallery")
def gallery():
    return render_template("gallery.html")


@app.route("/ppt_generator")
def ppt_generator_page():
    if "uid" not in session:
        return redirect(url_for("login"))
    return render_template("ppt_generator.html", session=session)


@app.route(
    "/generate-ppt",
    methods=["POST"]
)
def generate_ppt():

    try:

        if "file" not in request.files:

            return jsonify({
                "success": False,
                "error": "No file uploaded"
            })

        file = request.files["file"]

        slides = int(
            request.form.get(
                "slides",
                20
            )
        )

        filename = secure_filename(
            file.filename
        )

        filepath = os.path.join(
            NOTES_UPLOAD_FOLDER,
            filename
        )

        file.save(filepath)

        ppt_filename = f"ppt_{int(time.time())}.pptx"

        save_path, slides_data = create_ppt(
            filepath,
            ppt_filename,
            slides
        )

        conn = get_db()
        cur = conn.cursor()

        cur.execute("""
        INSERT INTO generated_files
        (
            file_name,
            generated_file,
            file_type,
            slide_json,
            created_at
        )
        VALUES (?, ?, ?, ?, ?)
        """,
        (
            filename,
            ppt_filename,
            "pptx",
            json.dumps(slides_data),
            datetime.now()
        ))

        conn.commit()
        conn.close()

        return jsonify({
            "success": True,
            "download_url": url_for(
                "download_generated",
                filename=ppt_filename
            )
        })

    except Exception as e:

        traceback.print_exc()

        return jsonify({
            "success": False,
            "error": str(e)
        })

@app.route(
    "/generate-notes",
    methods=["POST"]
)
def generate_notes():

    try:

        if "file" not in request.files:

            return jsonify({
                "success": False,
                "error": "No file uploaded"
            })

        file = request.files["file"]

        sections = int(
            request.form.get(
                "sections",
                10
            )
        )

        filename = secure_filename(
            file.filename
        )

        filepath = os.path.join(
            NOTES_UPLOAD_FOLDER,
            filename
        )

        file.save(filepath)

        notes_filename = f"notes_{int(time.time())}.pdf"

        save_path, notes_data = create_notes_pdf(
            filepath,
            notes_filename,
            sections
        )

        conn = get_db()
        cur = conn.cursor()

        cur.execute("""
        INSERT INTO generated_files
        (
            file_name,
            generated_file,
            file_type,
            created_at
        )
        VALUES (?, ?, ?, ?)
        """,
        (
            filename,
            notes_filename,
            "pdf",
            datetime.now()
        ))

        conn.commit()
        conn.close()

        return jsonify({
            "success": True,
            "download_url": url_for(
                "download_generated",
                filename=notes_filename
            )
        })

    except Exception as e:

        traceback.print_exc()

        return jsonify({
            "success": False,
            "error": str(e)
        })

@app.route(
    "/download/generated/<filename>"
)
def download_generated(filename):

    return send_file(
        os.path.join(
            GENERATED_FOLDER,
            filename
        ),
        as_attachment=False
    )


@app.route("/previous-year-papers")
def previous_year_papers():
    db = get_db()
    papers_rows = db.execute("SELECT * FROM previous_year_papers ORDER BY year DESC, semester ASC").fetchall()
    papers = []
    for row in papers_rows:
        paper = dict(row)
        photos = db.execute("SELECT image_path FROM paper_photos WHERE paper_id = ?", [paper.get('id')]).fetchall()
        paper['photos'] = [p['image_path'] for p in photos]
        papers.append(paper)
    db.close()
    return render_template("previous_year_papers.html", session=session, papers=papers)

@app.route("/add-paper", methods=["POST"])
def add_paper():
    if not session.get("uid"):
        return redirect(url_for('login'))
    
    subject = request.form.get("subject", "").strip()
    branch = request.form.get("branch", "").strip()
    year = request.form.get("year")
    semester = request.form.get("semester")
    photo_limit = request.form.get("photo_limit", 1)
    
    db = get_db()
    # Check for existing paper to prevent redundancy/duplication
    existing = db.execute(
        "SELECT id FROM previous_year_papers WHERE LOWER(subject) = LOWER(?) AND LOWER(branch) = LOWER(?) AND year = ? AND semester = ?",
        [subject, branch, year, semester]
    ).fetchone()
    
    if existing:
        db.close()
        return "Error: This paper record already exists for the selected year and semester.", 400

    db.execute(
        "INSERT INTO previous_year_papers (subject, branch, year, semester, photo_limit) VALUES (?, ?, ?, ?, ?)",
        [subject, branch, year, semester, photo_limit]
    )
    db.commit()
    db.close()
    return redirect(url_for('previous_year_papers'))

@app.route("/upload-paper-camera", methods=["POST"])
def upload_paper_camera():
    if not session.get("uid"):
        return jsonify({"success": False, "message": "Login required"})
    
    paper_id = request.form.get("paper_id")
    image_data = request.files.get("image")
    
    if image_data and paper_id:
        db = get_db()
        # Check limit
        paper = db.execute("SELECT photo_limit FROM previous_year_papers WHERE id = ?", [paper_id]).fetchone()
        current_count = db.execute("SELECT COUNT(*) as count FROM paper_photos WHERE paper_id = ?", [paper_id]).fetchone()
        
        if paper and current_count['count'] >= paper['photo_limit']:
            db.close()
            return jsonify({"success": False, "message": f"Maximum photo limit ({paper['photo_limit']}) reached for this paper."})

        filename = f"paper_{paper_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.jpg"
        filepath = os.path.join(app.static_folder, "images", filename)
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        image_data.save(filepath)
        
        db.execute(
            "INSERT INTO paper_photos (paper_id, image_path) VALUES (?, ?)",
            [paper_id, filename]
        )
        db.commit()
        db.close()
        return jsonify({"success": True, "filename": filename})
    
    return jsonify({"success": False})

# =====================================================================
# LOGIN (UID‑BASED)
# =====================================================================

@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "GET":
        message = request.args.get("message")
        return render_template("login.html", error=None, message=message)

    uid = request.form.get("uid") or request.form.get("username") or request.form.get("identifier")
    password = request.form.get("password")
    role = request.form.get("role") or request.form.get("mode")

    is_ajax = request.is_json or request.headers.get("X-Requested-With") == "XMLHttpRequest"
    if is_ajax and request.is_json:
        data = request.get_json() or {}
        uid = data.get("uid") or data.get("username") or uid
        password = data.get("password") or password
        role = data.get("role") or data.get("mode") or role

    uid = str(uid).strip() if uid else None
    role = str(role).strip().lower() if role else None

    uid_int = int(uid) if uid and uid.isdigit() else None

    db = get_db()
    row = db.execute(
        """
        SELECT u.uid, u.username, u.password_hash, u.role, u.name,
               u.email, u.phone, u.is_hod, s.roll
        FROM users u
        LEFT JOIN students s ON u.uid = s.uid
        WHERE u.uid = ? OR u.uid = ? OR u.username = ? OR u.email = ? OR s.roll = ?
        """,
        [uid, uid_int, uid, uid, uid]
    ).fetchone()

    if not row:
        db.close()
        msg = "Invalid UID or ID not found. Please check your credentials."
        if is_ajax:
            return jsonify({"success": False, "message": msg})
        return render_template("login.html", error=msg)

    db_role = row["role"]
    user_is_hod = int(row["is_hod"] or 0) == 1
    db_uid = row["uid"]
    db_hash = row["password_hash"]
    full_name = row["name"]

    entered_hash = hashlib.sha256((password or "").encode()).hexdigest()
    if db_hash != entered_hash:
        db.close()
        if is_ajax:
            return jsonify({"success": False, "message": "Invalid password"})
        return render_template("login.html", error="Invalid password")

    if role == "student" and db_role != "student":
        db.close()
        msg = "This ID is registered as Faculty/HOD. Please select the correct role."
        if is_ajax:
            return jsonify({"success": False, "message": msg})
        return render_template("login.html", error=msg)

    if role == "faculty" and db_role != "faculty":
        db.close()
        msg = "This ID is registered as a Student. Please select the correct role."
        if is_ajax:
            return jsonify({"success": False, "message": msg})
        return render_template("login.html", error=msg)

    db.close()

    session["uid"] = db_uid
    session["role"] = db_role
    session["name"] = full_name
    session["is_hod"] = user_is_hod

    if is_ajax:
        return jsonify({"success": True, "message": "Login successful!", "role": db_role})

    if db_role == "student":
        return redirect(url_for("student_dashboard"))
    return redirect(url_for("faculty_dashboard"))

# =====================================================================
# FORGOT PASSWORD
# =====================================================================

@app.route("/forgot-password", methods=["GET", "POST"])
def forgot_password():
    if request.method == "GET":
        return render_template("forget_password.html", error=None, step="request")

    action = request.form.get("action", "request")
    if action == "request":
        identifier = clean_text(request.form.get("identifier", ""))
        if not identifier:
            return render_template("forget_password.html", error="Missing identifier", step="request")

        db = get_db()
        row = db.execute(
            """
            SELECT u.uid, u.email, u.phone, s.roll
            FROM users u
            LEFT JOIN students s ON u.uid = s.uid
            WHERE u.uid = ? OR u.email = ? OR u.phone = ? OR s.roll = ?
            """,
            [identifier, identifier, identifier, identifier]
        ).fetchone()
        db.close()

        if not row:
            return render_template(
                "forget_password.html",
                error="No user found with that UID, email, phone, or roll number.",
                step="request"
            )

        user_email = row["email"] if isinstance(row, sqlite3.Row) else row[1]
        if not user_email:
            return render_template(
                "forget_password.html",
                error="No email address is available for this account.",
                step="request"
            )

        return render_template(
            "forget_password.html",
            step="reset",
            uid=row["uid"],
            email=user_email,
            message=f"A reset email would be sent to {user_email}. Enter your new password below."
        )

    if action == "reset":
        uid = request.form.get("uid")
        email = clean_text(request.form.get("email", ""))
        password = request.form.get("password", "")
        password_confirm = request.form.get("password_confirm", "")

        if not all([uid, email, password, password_confirm]):
            return render_template(
                "forget_password.html",
                error="All fields are required.",
                step="reset",
                uid=uid,
                email=email
            )

        if password != password_confirm:
            return render_template(
                "forget_password.html",
                error="Passwords do not match.",
                step="reset",
                uid=uid,
                email=email
            )

        db = get_db()
        row = db.execute(
            "SELECT email FROM users WHERE uid = ?",
            [uid]
        ).fetchone()
        if not row or (row["email"] if isinstance(row, sqlite3.Row) else row[0]) != email:
            db.close()
            return render_template(
                "forget_password.html",
                error="Email verification failed.",
                step="request"
            )

        password_hash = hashlib.sha256(password.encode()).hexdigest()
        db.execute(
            "UPDATE users SET password_hash = ? WHERE uid = ?",
            [password_hash, uid]
        )
        db.commit()
        db.close()

        return render_template(
            "forget_password.html",
            message="Password reset successfully. You can now login with your new password.",
            step="done"
        )


# =====================================================================
# STUDENT DASHBOARD
# =====================================================================

@app.route("/student-dashboard")
def student_dashboard():
    uid = session.get("uid")
    if not uid:
        return redirect(url_for("login"))

    if session.get("role") != "student":
        return redirect(url_for("login"))

    db = get_db()

    row = db.execute(
        """
        SELECT
            s.roll, s.name, s.department, s.cgpa, s.sgpa, s.attendance
        FROM
            students s
        WHERE
            s.uid = ?
        """,
        [uid]
    ).fetchone()
    if not row:
        db.close()
        return render_template(
            "student_dashboard.html",
            roll="N/A",
            name="Unknown Student",
            department="N/A",
            cgpa=0,
            sgpa=0,
            attendance=0,
            resumes=[],
            error="Student profile not found. Please contact the administrator.",
        ), 404

    roll, name, department, cgpa, sgpa, attendance = row

    resumes = db.execute(
        """
        SELECT
            r.id, r.uid, r.title, r.created_at
        FROM
            resumes r
        WHERE
            r.uid = ?
        ORDER BY
            r.created_at DESC
        """,
        [uid]
    ).fetchall()
    db.close()

    return render_template(
        "student_dashboard.html",
        roll=roll, name=name, department=department,
        cgpa=cgpa, sgpa=sgpa, attendance=attendance,
        resumes=resumes,
        session=session  # ← ADD THIS
    )

@app.route("/resume_form")
def resume_form():
    if "uid" not in session:
        return redirect(url_for("login"))
    return render_template("resume_form.html", session=session)

@app.route("/generate", methods=["POST"])
def generate():
    form_data = request.form.to_dict(flat=False)

    # Flatten single-item lists but keep lists for multi-input fields (like projects)
    for key, value in form_data.items():
        if len(value) == 1:
            form_data[key] = value[0]
    
    template_choice = form_data.get("template", "classic")
    
    # Handle profile photo upload
    photo = request.files.get("profile_photo")
    if photo and photo.filename:
        out_dir = os.path.join(BASE_DIR, "output")
        os.makedirs(out_dir, exist_ok=True)
        # Use session UID to keep temp files organized
        uid_prefix = session.get('uid', 'anon')
        photo_filename = f"temp_{uid_prefix}_{photo.filename}"
        photo_path = os.path.join(out_dir, photo_filename)
        photo.save(photo_path)
        form_data['photo_path'] = photo_path

    pdf_content = generate_resume_pdf(template_choice, form_data)
    
    # Save or update to database
    uid = session.get("uid")
    if uid:
        db = get_db()
        # Check if a resume for this user and template already exists
        existing_resume = db.execute(
            "SELECT id FROM resumes WHERE uid = ? AND title LIKE ?",
            [uid, f"Resume ({template_choice.capitalize()})%"]
        ).fetchone()

        if existing_resume:
            db.execute("UPDATE resumes SET pdf_content = ?, created_at = ? WHERE id = ?",
                       [pdf_content, datetime.now().strftime('%Y-%m-%d %H:%M:%S'), existing_resume["id"]])
        else:
            db.execute(
                "INSERT INTO resumes (uid, student_id, title, pdf_content, created_at) VALUES (?, ?, ?, ?, ?)",
                [uid, uid, f"Resume ({template_choice.capitalize()})", pdf_content, datetime.now().strftime('%Y-%m-%d %H:%M:%S')]
        )
        db.commit()
        db.close()

    filename = f"{clean_text(form_data.get('name') or 'resume').replace(' ', '_')}_{template_choice}.pdf"
    
    # Clean up temp photo if it exists
    if 'photo_path' in form_data and os.path.exists(form_data['photo_path']):
        try: os.remove(form_data['photo_path'])
        except: pass

    return send_file(
        io.BytesIO(pdf_content), 
        as_attachment=False, 
        download_name=filename, 
        mimetype="application/pdf"
    )

# =====================================================================
# FACULTY DEPARTMENT / RESUMES PAGES
# =====================================================================

@app.route("/faculty/<dept>")
def faculty_department(dept):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))

    db = get_db()
    students = db.execute(
        """
        SELECT s.*, u.email, u.phone 
        FROM students s 
        JOIN users u ON s.uid = u.uid 
        WHERE s.department = ? 
        ORDER BY s.roll
        """,
        [dept]
    ).fetchall()
    resumes = db.execute(
        """
        SELECT r.id, r.uid, r.title, r.created_at,
               s.roll, s.name
        FROM resumes r
        JOIN students s ON r.uid = s.uid
        WHERE s.department = ?
        ORDER BY s.roll
        """,
        [dept]
    ).fetchall()

    # Group resumes by student to avoid repeated files in the table
    from collections import defaultdict
    grouped = defaultdict(list)
    for r in resumes:
        grouped[r['uid']].append(dict(r))

    processed_resumes = []
    for uid, files in grouped.items():
        processed_resumes.append({
            'uid': uid,
            'roll': files[0]['roll'],
            'name': files[0]['name'],
            'latest_date': files[0]['created_at'],
            'files': files
        })
    processed_resumes.sort(key=lambda x: x['roll'])

    db.close()

    return render_template(
        "faculty-department.html",
        department=dept,
        students=students,
        resumes=processed_resumes,
        session=session
    )

@app.route("/debug_dashboard")
def debug_dashboard():
    uid = session.get("uid")
    if not uid or session.get("role") != "faculty":
        return "Login as faculty first"
    
    try:
        db = get_db()
        students_ncsml = db.execute("SELECT * FROM students WHERE department = 'NCSML'").fetchall()
        print("First NCSML student keys:", students_ncsml[0].keys() if students_ncsml else "NO DATA")
        print("First student type:", type(students_ncsml[0]) if students_ncsml else "EMPTY")
        db.close()
        return f"NCSML students: {len(students_ncsml)}<br>Type: {type(students_ncsml[0]) if students_ncsml else 'empty'}"
    except Exception as e:
        import traceback
        return f"ERROR: {str(e)}<br>{traceback.format_exc()}"

# =====================================================================
# EDIT STUDENT DETAILS - FACULTY DASHBOARD
# =====================================================================
@app.route("/faculty/edit_student/<uid>", methods=["GET", "POST"])
def faculty_edit_student(uid):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))

    db = get_db()

    if request.method == "GET":
        row = db.execute(
            "SELECT s.*, u.email, u.phone FROM students s JOIN users u ON s.uid = u.uid WHERE s.uid = ?",
            [uid]
        ).fetchone()
        
        # Convert Row to dict so Jinja2 can use dot notation (e.g. student.name)
        student = dict(row) if row else None

        if not student:
            db.close()
            return "Student not found", 404

        db.close()
        return render_template("faculty-edit-student.html", student=student)

    # POST
    department = request.form.get("department")
    roll = request.form.get("roll")
    name = request.form.get("name")
    email = request.form.get("email")
    phone = request.form.get("phone")
    cgpa = float(request.form.get("cgpa", 0.0) or 0.0)
    sgpa = float(request.form.get("sgpa", 0.0) or 0.0)
    attendance = int(request.form.get("attendance", 0) or 0)
    

    if not all([department, roll, name]):
        db.close()
        return "Missing required fields", 400

    try:
        # Update users table to keep name, email and phone in sync
        db.execute(
            """
            UPDATE users
            SET name = ?, email = ?, phone = ?
            WHERE uid = ?
            """,
            [name, email, phone, uid]
        )

        # Update students table
        db.execute(
            """
            UPDATE students
            SET department = ?, roll = ?, name = ?, cgpa = ?, sgpa = ?, attendance = ?
            WHERE uid = ?
            """,
            [department, roll, name, cgpa, sgpa, attendance, uid]
        )
        db.commit()
    except Exception as e:
        db.rollback()
        return f"Error updating student: {str(e)}", 500
    finally:
        db.close()

    return redirect(url_for("faculty_dashboard"))


# =====================================================================
# FACULTY ADD STUDENT PAGE
# =====================================================================
@app.route("/faculty/add_student", methods=["GET"])
def faculty_add_student_form():

    if session.get("role") != "faculty":
        return redirect(url_for("login"))

    return render_template("faculty_add_student.html")


# =====================================================================
# FACULTY ADD STUDENT SAVE
# =====================================================================
@app.route("/faculty/add_student", methods=["POST"])
def faculty_add_student():

    if session.get("role") != "faculty":
        return redirect(url_for("login"))

    username = request.form.get("username")
    password = request.form.get("password")
    department = request.form.get("department")
    roll = request.form.get("roll")
    name = request.form.get("name")

    email = request.form.get("email")
    phone = request.form.get("phone")

    cgpa = float(request.form.get("cgpa", 0.0) or 0.0)
    sgpa = float(request.form.get("sgpa", 0.0) or 0.0)
    attendance = int(request.form.get("attendance", 0) or 0)

    if not all([username, password, department, roll, name]):
        return "Missing required fields", 400

    db = get_db()

    try:
        h = hashlib.sha256(password.encode()).hexdigest()

        db.execute(
            """
            INSERT INTO users
            (username, password_hash, name, role, email, phone)
            VALUES (?, ?, ?, 'student', ?, ?)
            """,
            [username, h, name, email, phone]
        )

        db.commit()

        user_row = db.execute(
            "SELECT uid FROM users WHERE username = ?",
            [username]
        ).fetchone()

        uid = user_row[0]

        db.execute(
            """
            INSERT INTO students
            (uid, department, roll, name, cgpa, sgpa, attendance)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            [uid, department, roll, name, cgpa, sgpa, attendance]
        )

        db.commit()

    except Exception as e:
        db.rollback()
        db.close()
        return str(e)


    return redirect(url_for("faculty_dashboard"))
# =====================================================================
# FACULTY DELETE STUDENT SAVE
# =====================================================================

@app.route("/faculty/delete_student/<string:roll>", methods=["POST"])
def faculty_delete_student(roll):

    if session.get("role") != "faculty":
        return redirect(url_for("login"))

    db = get_db()

    row = db.execute(
        "SELECT uid FROM students WHERE roll = ?",
        [roll]
    ).fetchone()

    if row:
        uid = row[0]

        db.execute(
            "DELETE FROM resumes WHERE uid = ?",
            [uid]
        )

        db.execute(
            "DELETE FROM students WHERE roll = ?",
            [roll]
        )

        db.execute(
            "DELETE FROM users WHERE uid = ?",
            [uid]
        )

        db.commit()

    db.close()
    return redirect(url_for("faculty_dashboard"))

@app.route("/notes")
def notes():
    if not session.get("uid"):
        return redirect(url_for("login"))
    db = get_db()
    notes_rows = db.execute("""
        SELECT n.*, u.name as uploader_name 
        FROM notes n 
        JOIN users u ON n.uploaded_by_uid = u.uid 
        ORDER BY n.created_at DESC
    """).fetchall()
    
    notes_list = []
    for row in notes_rows:
        note = dict(row)
        files = db.execute("SELECT * FROM note_files WHERE note_id = ?", [note['id']]).fetchall()
        note['files'] = files
        notes_list.append(note)
    db.close()
    return render_template("notes.html", notes=notes_list, session=session)

@app.route("/faculty/add-notes", methods=["GET", "POST"])
def faculty_add_notes():
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    
    if request.method == "POST":
        subject = request.form.get("subject")
        semester = request.form.get("semester")
        department = request.form.get("department")
        year = request.form.get("year")
        files = request.files.getlist("note_files")
        
        db = get_db()
        cursor = db.cursor()
        cursor.execute(
            "INSERT INTO notes (subject, semester, department, year, uploaded_by_uid) VALUES (?, ?, ?, ?, ?)",
            [subject, semester, department, year, session["uid"]]
        )
        note_id = cursor.lastrowid
        
        for file in files:
            if file and allowed_notes_file(file.filename):
                filename = secure_filename(file.filename)
                unique_filename = f"{int(time.time())}_{filename}"
                file_path = os.path.join(NOTES_UPLOAD_FOLDER, unique_filename)
                file.save(file_path)
                file_type = filename.rsplit('.', 1)[1].lower()
                cursor.execute(
                    "INSERT INTO note_files (note_id, filename, file_path, file_type) VALUES (?, ?, ?, ?)",
                    [note_id, unique_filename, unique_filename, file_type]
                )
        db.commit()
        db.close()
        return redirect(url_for("notes"))
    return render_template("faculty_add_notes.html")

@app.route("/faculty_dashboard")
def faculty_dashboard():
    try:
        uid = session.get("uid")
        if not uid or session.get("role") != "faculty":
            return redirect(url_for("login"))

        dept = request.args.get("dept", None) 

        db = get_db()
        
        def get_students_for(department):
            return db.execute(
                "SELECT s.*, u.email, u.phone FROM students s JOIN users u ON s.uid = u.uid WHERE s.department = ? ORDER BY s.roll",
                [department]
            ).fetchall()

        def get_resumes_for(department):
            return db.execute(
                "SELECT r.id, r.uid, r.title, r.created_at, s.roll, s.name FROM resumes r JOIN students s ON r.uid = s.uid WHERE s.department = ? ORDER BY r.created_at DESC",
                [department]
            ).fetchall()

        context = {
            "students_acsml": get_students_for("ACSML"),
            "students_ncsml": get_students_for("NCSML"),
            "students_dcsml": get_students_for("DCSML"),
            "resumes_acsml": get_resumes_for("ACSML"),
            "resumes_ncsml": get_resumes_for("NCSML"),
            "resumes_dcsml": get_resumes_for("DCSML"),
            "session": session,
            "dept": dept
        }
        db.close()
        return render_template("faculty_dashboard.html", **context)

    except Exception as e:
        import traceback
        return f"<pre>{traceback.format_exc()}</pre>"

# =====================================================================
# RESUME TEMPLATE GENERATORS (5 Styles)
# =====================================================================
@app.route("/signup", methods=["GET", "POST"])
def signup():
    if request.method == "POST":
        # Support both standard form submission and AJAX (JSON) requests
        if request.is_json:
            data = request.get_json()
            name = data.get("name")
            uid = data.get("uid")
            email = data.get("email")
            password = data.get("password")
            role = data.get("role")
            phone = data.get("phone", "")
        else:
            name = request.form.get("name")
            uid = request.form.get("uid")
            email = request.form.get("email")
            password = request.form.get("password")
            role = request.form.get("role")
            phone = request.form.get("phone", "")

        # Handle HOD permission during account creation
        is_hod_val = 1 if role == "hod" else 0
        db_role_val = "faculty" if role == "hod" else role

        h = hashlib.sha256(password.encode()).hexdigest()
        db = get_db()
        try:
            db.execute(
                "INSERT INTO users (uid, username, password_hash, name, role, email, phone, is_hod) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                [uid, uid, h, name, db_role_val, email, phone, is_hod_val]
            )
            if db_role_val == 'student':
                db.execute(
                    "INSERT INTO students (uid, name, department, roll, cgpa, sgpa, attendance) VALUES (?, ?, 'ACSML', ?, 0.0, 0.0, 0)",
                    [uid, name, uid]
                )
            db.commit()
            if request.is_json:
                return jsonify({"success": True, "message": "Account created successfully!"})
            return redirect(url_for('login', message="Account created successfully!"))
        except sqlite3.IntegrityError:
            if request.is_json:
                return jsonify({"success": False, "message": "UID, Username or Email already exists."})
            return render_template("signup.html", error="UID or Email already exists.")
        finally:
            db.close()
    return render_template("signup.html")


def make_styles(template_type):
    styles = getSampleStyleSheet()
    body_font = "Times-Roman"
    bold_font = "Times-Bold"

    palette = {
        "classic": ("#1e3a5f", "#51606f", "#d7dee6"),
        "modern": ("#0f766e", "#5b6472", "#dbe4ea"),
        "minimal": ("#111827", "#6b7280", "#e5e7eb"),
        "academic": ("#334155", "#64748b", "#cbd5e1"),
        "creative": ("#8b5e3c", "#647b68", "#e8dfd4"),
    }
    primary, muted, line = palette.get(template_type, palette["classic"])
    primary = colors.HexColor(primary)
    muted = colors.HexColor(muted)
    line = colors.HexColor(line)

    return {
        "title": ParagraphStyle("title", parent=styles["Normal"], fontName=bold_font, fontSize=22, leading=22*1.15, textColor=primary, alignment=TA_LEFT, spaceAfter=3),
        "subtitle": ParagraphStyle("subtitle", parent=styles["Normal"], fontName=body_font, fontSize=10, leading=10*1.15, textColor=muted, spaceAfter=7),
        "section": ParagraphStyle("section", parent=styles["Heading2"], fontName=bold_font, fontSize=11.1, leading=11.1*1.15, textColor=primary, spaceBefore=6, spaceAfter=5),
        "body": ParagraphStyle("body", parent=styles["BodyText"], fontName=body_font, fontSize=10, leading=10*1.15, textColor=colors.black),
        "small": ParagraphStyle("small", parent=styles["BodyText"], fontName=body_font, fontSize=9, leading=9*1.15, textColor=muted),
        "center_title": ParagraphStyle("center_title", parent=styles["Normal"], fontName=bold_font, fontSize=22, leading=26, textColor=primary, alignment=TA_CENTER, spaceAfter=2),
        "center_sub": ParagraphStyle("center_sub", parent=styles["Normal"], fontName=body_font, fontSize=10, leading=12, textColor=muted, alignment=TA_CENTER),
        "right_title": ParagraphStyle("right_title", parent=styles["Normal"], fontName=bold_font, fontSize=18, leading=22, textColor=primary, alignment=TA_RIGHT, spaceAfter=2),
        "justify": ParagraphStyle("justify", parent=styles["BodyText"], fontName=body_font, fontSize=9.2, leading=12.5, textColor=colors.black, alignment=TA_JUSTIFY),
        "line": line,
        "primary": primary,
        "muted": muted,
    }

def page_header(canvas, doc, form_data, template_type):
    canvas.saveState()
    w, h = A4

    if template_type == "modern":
        canvas.setFillColor(colors.HexColor("#ecfdf5"))
        canvas.rect(0, h - 48, w, 48, stroke=0, fill=1)
        canvas.setStrokeColor(colors.HexColor("#0f766e"))
        canvas.setLineWidth(1.2)
        canvas.line(doc.leftMargin, h - 24, w - doc.rightMargin, h - 24)
        canvas.setFont("Helvetica", 9)
        canvas.setFillColor(colors.HexColor("#0f766e"))
        canvas.drawRightString(w - doc.rightMargin, h - 16, clean_text(form_data.get("field_of_study")))

    elif template_type == "academic":
        canvas.setStrokeColor(colors.HexColor("#94a3b8"))
        canvas.setLineWidth(1)
        canvas.line(doc.leftMargin, h - 22, w - doc.rightMargin, h - 22)
        canvas.setFont("Helvetica", 8.5)
        canvas.setFillColor(colors.HexColor("#64748b"))
        canvas.drawString(doc.leftMargin, h - 16, clean_text(form_data.get("field_of_study")))

    elif template_type == "creative":
        canvas.setFillColor(colors.HexColor("#f7f3ee"))
        canvas.rect(0, h - 52, w, 52, stroke=0, fill=1)
        canvas.setStrokeColor(colors.HexColor("#8b5e3c"))
        canvas.line(doc.leftMargin, h - 26, w - doc.rightMargin, h - 26)

    elif template_type == "minimal":
        canvas.setStrokeColor(colors.HexColor("#111827"))
        canvas.line(25 * mm, h - 24, w - 25 * mm, h - 24)

    else:
        canvas.setStrokeColor(colors.HexColor("#d7dee6"))
        canvas.line(doc.leftMargin, h - 26, w - doc.rightMargin, h - 26)
        canvas.setFont("Helvetica", 9)
        canvas.setFillColor(colors.HexColor("#51606f"))
        canvas.drawString(doc.leftMargin, h - 18, clean_text(form_data.get("field_of_study")))

    canvas.restoreState()

def box(title, text, styles, fill=None):
    tbl = Table([[Paragraph(f"<b>{title}</b>", styles["body"])],
                 [Paragraph(text, styles["small"])]],
                colWidths=[None])
    ts = [
        ("BOX", (0, 0), (-1, -1), 0.8, styles["line"]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]
    if fill:
        ts.append(("BACKGROUND", (0, 0), (-1, -1), fill))
    tbl.setStyle(TableStyle(ts))
    return tbl

def common_sections(form_data, styles):
    story = []

    story.append(Paragraph("CAREER OBJECTIVE", styles["section"]))
    story.append(Paragraph(
        enhance_objective(form_data.get("objective"), form_data.get("field_of_study")),
        styles["body"]
    ))
    story.append(Spacer(1, 5))

    story.append(Paragraph("ACADEMIC QUALIFICATION", styles["section"]))

    if form_data.get("ssc_school"):
        ssc = (
            f"<b>SSC (10th Standard)</b>: {clean_text(form_data.get('ssc_school'))}<br/>"
            f"Score: {clean_text(form_data.get('ssc_percentage'))}"
        )
        story.append(Paragraph(ssc, styles["body"]))
        story.append(Spacer(1, 3))

    if form_data.get("inter_college"):
        inter = (
            f"<b>Intermediate (12th Standard)</b>: {clean_text(form_data.get('inter_college'))}<br/>"
            f"Score: {clean_text(form_data.get('inter_percentage'))}"
        )
        story.append(Paragraph(inter, styles["body"]))
        story.append(Spacer(1, 3))

    if form_data.get("ug_college"):
        ug = (
            f"<b>UG</b>: {clean_text(form_data.get('ug_degree'))} in {clean_text(form_data.get('ug_branch'))}<br/>"
            f"{clean_text(form_data.get('ug_college'))}<br/>"
            f"Score: {clean_text(form_data.get('ug_score'))}"
        )
        story.append(Paragraph(ug, styles["body"]))
        story.append(Spacer(1, 3))

    if form_data.get("pg_college"):
        pg = (
            f"<b>PG</b>: {clean_text(form_data.get('pg_degree'))} in {clean_text(form_data.get('pg_branch'))}<br/>"
            f"{clean_text(form_data.get('pg_college'))}<br/>"
            f"Score: {clean_text(form_data.get('pg_score'))}"
        )
        story.append(Paragraph(pg, styles["body"]))
        story.append(Spacer(1, 3))

    # This block was duplicating UG info. Removed to keep sections distinct.

    # Process skills from form_data (skill_name and skill_level)
    skill_names = form_data.get("skill_name", [])
    skill_levels = form_data.get("skill_level", [])

    processed_skills = []

    if isinstance(skill_names, str):
        skill_names = [skill_names]
        skill_levels = [skill_levels]

    for i in range(len(skill_names)):
        name = clean_text(skill_names[i])
        level = clean_text(skill_levels[i])

        if name:
            processed_skills.append(f"{name} - {level}")
    form_data['processed_technical_skills'] = " • ".join(processed_skills)
    return story
def parse_projects_with_bullets(form_data):
    print("RAW DATA:", form_data)

    titles = form_data.get("project_title", [])
    descriptions = form_data.get("project_description", [])

    print("TITLES =", titles)
    print("DESCRIPTIONS =", descriptions)

    if not isinstance(titles, list):
        titles = [titles]

    if not isinstance(descriptions, list):
        descriptions = [descriptions]

    projects = []

    for title, desc in zip(titles, descriptions):
        projects.append({
            "title": str(title),
            "bullets": [str(desc)]
        })

    return projects
def generate_resume_pdf(template_type, form_data):
    buffer = io.BytesIO()
    template_type = (template_type or "classic").lower()
    if template_type not in ["classic", "modern", "minimal", "academic", "creative"]:
        template_type = "classic"

    styles = make_styles(template_type)

    def on_page(canvas, doc):
        page_header(canvas, doc, form_data, template_type)
        canvas.saveState()
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(colors.HexColor("#6b7280"))
        canvas.drawRightString(A4[0] - doc.rightMargin, 8 * mm, f"Page {doc.page}")
        canvas.restoreState()

    if template_type == "modern":
        doc = BaseDocTemplate(
            buffer,
            pagesize=A4,
            leftMargin=12*mm,
            rightMargin=12*mm,
            topMargin=54*mm,
            bottomMargin=15*mm
        )
        sidebar_w = 60 * mm
        gap = 6 * mm
        main_w = A4[0] - doc.leftMargin - doc.rightMargin - sidebar_w - gap
        frame_h = A4[1] - doc.topMargin - doc.bottomMargin

        left_frame = Frame(doc.leftMargin, doc.bottomMargin, sidebar_w, frame_h, id="sidebar", showBoundary=0)
        right_frame = Frame(doc.leftMargin + sidebar_w + gap, doc.bottomMargin, main_w, frame_h, id="main", showBoundary=0)
        doc.addPageTemplates([PageTemplate(id="modern", frames=[left_frame, right_frame], onPage=on_page)])

        story = []
        story.append(Paragraph("PROFILE", styles["section"]))
        story.append(Paragraph(enhance_objective(form_data.get("objective"), form_data.get("field_of_study")), styles["small"]))
        story.append(Spacer(1, 4))

        story.append(Paragraph("CONTACT", styles["section"]))
        contact_lines = [
            f"Email: {clean_text(form_data.get('email'))}",
            f"Phone: {clean_text(form_data.get('phone'))}"
        ]
        if form_data.get('linkedin'): contact_lines.append(f"LinkedIn: {clean_text(form_data.get('linkedin'))}")
        if form_data.get('github'): contact_lines.append(f"GitHub: {clean_text(form_data.get('github'))}")
        if form_data.get('website'): contact_lines.append(f"Website: {clean_text(form_data.get('website'))}")
        location_parts = [clean_text(form_data.get('city')), clean_text(form_data.get('state')), clean_text(form_data.get('country'))]
        location_str = ", ".join(filter(None, location_parts))
        if location_str: contact_lines.append(f"Location: {location_str}")

        story.append(Paragraph("<br/>".join(contact_lines), styles["small"]))
        story.append(Spacer(1, 4))

        # Skills section for modern template
        story.append(Paragraph("SKILLS", styles["section"]))
        # Skills logic handled via processed_technical_skills in original, 
        # but we need to ensure it's calculated before this.
        skill_names = form_data.get("skill_name", [])
        skill_levels = form_data.get("skill_level", [])

        if isinstance(skill_names, str):
            skill_names = [skill_names]

        if isinstance(skill_levels, str):
            skill_levels = [skill_levels]

        processed_skills = []

        for i in range(len(skill_names)):
            name = clean_text(skill_names[i])

            level = skill_levels[i] if i < len(skill_levels) else ""

            if name:
                processed_skills.append(f"{name} - {level}")
        story.append(Paragraph(" • ".join(processed_skills) or "N/A", styles["small"]))
        story.append(Spacer(1, 4))

        story.append(Paragraph("LANGUAGES", styles["section"]))
        langs = "<br/>".join(f"• {clean_text(x)}" for x in split_items(form_data.get("languages")))
        story.append(Paragraph(langs or "N/A", styles["small"]))
        story.append(FrameBreak())

        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["subtitle"]))
        story.extend(common_sections(form_data, styles))

        projects = parse_projects_with_bullets(form_data)
        print("PARSED PROJECTS:", parse_projects_with_bullets(form_data))
        if projects:
            story.append(Paragraph("PROJECTS", styles["section"]))
            for i, p in enumerate(projects, 1):
                title = clean_text(p["title"]) or f"Project {i}"
                bullets = p["bullets"] or ["Project details not provided."]
                txt = f"<b>{title}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets)
                story.append(box(f"Project {i}", txt, styles, fill=colors.HexColor("#f0fdfa")))
                story.append(Spacer(1, 4))

        experiences = split_items(form_data.get("work_experience")) or split_items(form_data.get("internships"))
        if experiences:
            story.append(Paragraph("INTERNSHIPS", styles["section"]))
            for exp in experiences:
                bullets = describe_experience(exp)
                txt = f"<b>{clean_text(exp)}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets)
                story.append(Paragraph(txt, styles["body"]))
                story.append(Spacer(1, 4))

        certs = split_items(form_data.get("certifications"))
        if certs:
            story.append(Paragraph("CERTIFICATIONS", styles["section"]))
            for c in certs:
                desc = describe_certification(c)
                story.append(Paragraph(f"• <b>{clean_text(c)}</b>: {clean_text(desc)}", styles["body"]))
                story.append(Spacer(1, 2))

        doc.build(story)
        pdf_out = buffer.getvalue()
        buffer.close()
        return pdf_out

    doc = BaseDocTemplate(buffer, pagesize=A4, leftMargin=15*mm, rightMargin=15*mm, topMargin=18*mm, bottomMargin=15*mm)
    story = []

    if template_type == "classic":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        contact_parts = [
            clean_text(form_data.get('field_of_study')),
            clean_text(form_data.get('email')),
            clean_text(form_data.get('phone'))
        ]
        if form_data.get('linkedin'): contact_parts.append(f"LinkedIn: {clean_text(form_data.get('linkedin'))}")
        if form_data.get('github'): contact_parts.append(f"GitHub: {clean_text(form_data.get('github'))}")
        if form_data.get('website'): contact_parts.append(f"Website: {clean_text(form_data.get('website'))}")
        story.append(Paragraph(" | ".join(filter(None, contact_parts)), styles["subtitle"]))
        story.append(HRFlowable(width="100%", thickness=0.8, color=styles["line"]))
        story.append(Spacer(1, 5))

    elif template_type == "minimal":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["center_title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["center_sub"]))
        contact_parts = [clean_text(form_data.get('email')), clean_text(form_data.get('phone'))]
        if form_data.get('linkedin'): contact_parts.append(clean_text(form_data.get('linkedin')))
        if form_data.get('github'): contact_parts.append(clean_text(form_data.get('github')))
        story.append(Paragraph(" • ".join(filter(None, contact_parts)), styles["center_sub"]))
        story.append(Spacer(1, 5))
        story.append(HRFlowable(width="42%", thickness=1, color=styles["line"]))
        story.append(Spacer(1, 7))

    elif template_type == "academic":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["right_title"]))
        contact_info = f"{clean_text(form_data.get('email'))} | {clean_text(form_data.get('phone'))}"
        if form_data.get('linkedin'): contact_info += f" | LinkedIn: {clean_text(form_data.get('linkedin'))}"
        story.append(Paragraph(contact_info, styles["small"]))
        story.append(Spacer(1, 4))
        story.append(HRFlowable(width="100%", thickness=1, color=styles["line"]))
        story.append(Spacer(1, 6))

    elif template_type == "creative":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["subtitle"]))
        if form_data.get('website'): story.append(Paragraph(f"Portfolio: {clean_text(form_data.get('website'))}", styles["small"]))
        story.append(Spacer(1, 3))
        story.append(HRFlowable(width="100%", thickness=1.1, color=styles["line"]))
        story.append(Spacer(1, 7))

    story.extend(common_sections(form_data, styles))

    # Recalculate stars for other templates
    skill_names = form_data.get("skill_name", [])
    skill_levels = form_data.get("skill_level", [])

    if isinstance(skill_names, str):
        skill_names = [skill_names]

    if isinstance(skill_levels, str):
        skill_levels = [skill_levels]

    processed_skills = []

    for i in range(len(skill_names)):
        name = clean_text(skill_names[i])

        level = skill_levels[i] if i < len(skill_levels) else ""

        if name:
            processed_skills.append(f"{name} - {level}")

    if processed_skills:
        skill_str = " • ".join(processed_skills)
        story.append(Paragraph("TECHNICAL SKILLS", styles["section"]))
        story.append(Paragraph(skill_str, styles["body"]))
        story.append(Spacer(1, 5))

    # Languages section for other templates
    if form_data.get("languages"):
        story.append(Paragraph("LANGUAGES", styles["section"])) # This line was missing for non-modern templates
        if template_type == "minimal":
            data = [[Paragraph("Languages", styles["body"]), Paragraph(" • ".join(clean_text(x) for x in split_items(form_data.get("languages"))), styles["small"])]]
            t = Table(data, colWidths=[35*mm, None])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#f8fafc")),
                ("BOX", (0, 0), (-1, -1), 0.8, styles["line"]),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]))
            story.append(t)
        elif template_type == "academic":
            story.append(Paragraph(" • ".join(clean_text(x) for x in split_items(form_data.get("languages"))), styles["justify"]))
        else:
            story.append(Paragraph(" • ".join(clean_text(x) for x in split_items(form_data.get("languages"))), styles["body"]))
        story.append(Spacer(1, 5))

    projects = parse_projects_with_bullets(form_data)
    if projects:
        story.append(Paragraph("PROJECTS", styles["section"]))
        for i, p in enumerate(projects, 1):
            title = clean_text(p["title"]) or f"Project {i}"
            bullets = p["bullets"] or ["Project details not provided."]
            txt = f"<b>{title}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets)
            if template_type == "modern":
                story.append(box(f"Project {i}", txt, styles, fill=colors.HexColor("#f0fdfa")))
            elif template_type == "creative":
                story.append(box(f"Project {i}", txt, styles, fill=colors.HexColor("#fbf7f2")))
            elif template_type == "academic":
                story.append(Paragraph(txt, styles["justify"]))
            else:
                story.append(Paragraph(txt, styles["body"]))
            story.append(Spacer(1, 4))

    experiences = split_items(form_data.get("work_experience")) or split_items(form_data.get("internships"))
    if experiences:
        story.append(Paragraph("EXPERIENCE & INTERNSHIPS", styles["section"]))
        for exp in experiences:
            bullets = describe_experience(exp)
            txt = f"<b>{clean_text(exp)}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets)
            story.append(Paragraph(txt, styles["body"]))
            story.append(Spacer(1, 5))
        story.append(Spacer(1, 5))

    certs = split_items(form_data.get("certifications"))
    if certs:
        story.append(Paragraph("CERTIFICATIONS", styles["section"]))
        for c in certs:
            desc = describe_certification(c)
            story.append(Paragraph(f"• <b>{clean_text(c)}</b>: {clean_text(desc)}", styles["body"]))
            story.append(Spacer(1, 3))

    doc.addPageTemplates([PageTemplate(id=template_type, frames=[Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="main")], onPage=on_page)])
    doc.build(story)
    pdf_out = buffer.getvalue()
    buffer.close()
    return pdf_out

# =====================================================================
# DOWNLOAD 
# =====================================================================

@app.route("/download_resume/<int:resume_id>")
def download_resume(resume_id):
    db = get_db()
    row = db.execute("SELECT pdf_content, title FROM resumes WHERE id = ?", [resume_id]).fetchone()
    db.close()

    if not row:
        return "Resume not found in database", 404

    pdf_content = row["pdf_content"]
    filename = f"{row['title'].replace(' ', '_')}.pdf"
    return send_file(io.BytesIO(pdf_content), as_attachment=True, download_name=filename, mimetype="application/pdf")

# =====================================================================
# VIEW RESUME
# =====================================================================

@app.route("/view_resume/<int:resume_id>")
def view_resume(resume_id):
    if session.get("role") == "faculty":
        db = get_db()
        row = db.execute(
            "SELECT pdf_content FROM resumes WHERE id = ?",
            [resume_id]
        ).fetchone()
        db.close()
    else:
        db = get_db()
        row = db.execute(
            "SELECT pdf_content FROM resumes WHERE id = ? AND uid = ?",
            [resume_id, session.get("uid")]
        ).fetchone()
        db.close()

    if not row:
        return "Resume not found", 404

    pdf_content = row["pdf_content"]
    return send_file(io.BytesIO(pdf_content), mimetype="application/pdf")

# =====================================================================
#  FACULTY VIEW RESUME
# =====================================================================

@app.route("/faculty/resumes")
def faculty_resumes():
    if session.get("role") != "faculty":
        return redirect(url_for("login"))

    db = get_db()
    rows = db.execute("""
        SELECT r.id, r.uid, r.title, r.created_at, s.roll, s.name, s.department
        FROM resumes r
        JOIN students s ON r.uid = s.uid
        ORDER BY s.roll, r.created_at DESC
    """).fetchall()

    from collections import defaultdict
    grouped = defaultdict(list)
    for r in rows:
        grouped[r['uid']].append(dict(r))
    
    resumes = []
    for uid, st_resumes in grouped.items():
        resumes.append({
            'uid': uid,
            'roll': st_resumes[0]['roll'],
            'name': st_resumes[0]['name'],
            'department': st_resumes[0]['department'],
            'latest_date': st_resumes[0]['created_at'],
            'files': st_resumes
        })
    resumes.sort(key=lambda x: x['latest_date'], reverse=True)
    db.close()
    return render_template("faculty_resumes.html", resumes=resumes, session=session)

@app.route("/combined_resume/<uid>")
def combined_resume(uid):
    if not session.get("uid"):
        return redirect(url_for("login"))
    if session.get("role") != "faculty" and session.get("uid") != uid:
        return "Unauthorized", 403
    db = get_db()
    rows = db.execute("SELECT pdf_content FROM resumes WHERE uid = ? ORDER BY created_at ASC", [uid]).fetchall()
    db.close()
    if not rows:
        return "No resumes found", 404
    if len(rows) == 1:
        return send_file(io.BytesIO(rows[0]["pdf_content"]), mimetype="application/pdf")
    try:
        from pypdf import PdfWriter
        merger = PdfWriter()
        for row in rows:
            if row["pdf_content"]:
                merger.append(io.BytesIO(row["pdf_content"]))
        output = io.BytesIO()
        merger.write(output)
        output.seek(0)
        return send_file(output, as_attachment=True, download_name=f"Combined_Resumes_{uid}.pdf", mimetype="application/pdf")
    except ImportError:
        return send_file(io.BytesIO(rows[-1]["pdf_content"]), mimetype="application/pdf")


# =====================================================================
# DELETE RESUME
# =====================================================================

@app.route("/delete_resume", methods=["POST"])
def delete_resume():
    if "uid" not in session:
        return redirect(url_for("login"))

    resume_id_str = request.form.get("resume_id")
    if not resume_id_str:
        return "Invalid request", 400

    try:
        resume_id = int(resume_id_str)
    except ValueError:
        return "Invalid resume ID", 400

    uid = session["uid"]
    db = get_db()

    db.execute("DELETE FROM resumes WHERE id = ? AND uid = ?", [resume_id, uid])
    db.commit()
    db.close()

    return redirect(url_for("student_dashboard"))

# =====================================================================
# STUDENT DETAILS
# =====================================================================

@app.route("/student_details/<student_uid>")
def student_details(student_uid):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))

    db = get_db()
    db.row_factory = sqlite3.Row
    
    # Get student details with contact information
    student = db.execute(
        """
        SELECT s.uid, s.roll, s.name, s.department, s.cgpa, s.sgpa, s.attendance,
               u.email, u.phone
        FROM students s
        LEFT JOIN users u ON s.uid = u.uid
        WHERE s.uid = ?
        """,
        [student_uid]
    ).fetchone()
    
    if not student:
        db.close()
        return "Student not found", 404
    
    # Get student's resumes
    resumes = db.execute(
        "SELECT * FROM resumes WHERE uid = ? ORDER BY created_at DESC",
        [student_uid]
    ).fetchall()
    
    db.close()
    
    return render_template("student_details.html", student=student, resumes=resumes)
# =====================================================================
# NOTES
# ====================================================================

# =====================================================================
# TEST CASE
# =====================================================================
@app.route("/test")
def test():
    return render_template("test.html")


# =====================================================================
#  CHECK AUTH STATUS
# =====================================================================
@app.route("/check-auth")
def check_auth_endpoint():
    return jsonify({
        "logged_in": "uid" in session,
        "role": session.get("role")
    })


# =====================================================================
#  LOGOUT 
# =====================================================================
@app.context_processor
def inject_recent_resumes():
    if "uid" in session and session.get("role") == "student":
        try:
            db = get_db()
            rows = db.execute(
                "SELECT id, title FROM resumes WHERE uid = ? ORDER BY created_at DESC LIMIT 3",
                [session["uid"]]
            ).fetchall()
            recent = [dict(r) for r in rows]
            db.close()
            return dict(recent_resumes=recent)
        except:
            return dict(recent_resumes=[])
    return dict(recent_resumes=[])

@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))

@app.route("/generate_project_description", methods=["POST"])
def generate_project_description():
    data = request.get_json()
    title = data.get("title", "")
    features = data.get("features", "")

    if not client:
        logger.error("AI Generation Error: Gemini client is not initialized.")
        return jsonify({
            "success": False,
            "error": "AI service is currently unavailable. Please check server logs."
        }), 503

    if not GEMINI_API_KEY or GEMINI_API_KEY == "your_fallback_key_here" or not GEMINI_API_KEY.startswith("AIza"):
        logger.error(f"AI CONFIG ERROR: Invalid GEMINI_API_KEY. Prefix: {GEMINI_API_KEY[:4] if GEMINI_API_KEY else 'None'}...")
        return jsonify({
            "success": False,
            "error": "AI configuration error. A valid Google API key (starting with 'AIza') is required."
        }), 401

    prompt = f"""
    Project Title: {title}
    Features:
    {features}
    Write 4 professional ATS-friendly resume bullet points.
    Use action verbs.
    """

    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = client.models.generate_content(
                model="gemini-1.5-flash",
                contents=prompt,
                config={'http_options': {'timeout': 30}}
            )

            if response and response.text:
                return jsonify({
                    "success": True,
                    "description": response.text
                })
            else:
                raise ValueError("Empty response from AI model")

        except Exception as e:
            err_msg = str(e).lower()
            # Check for transient errors like 503 (Service Unavailable) or 429 (Rate Limit)
            if ("503" in err_msg or "429" in err_msg or "overloaded" in err_msg) and attempt < max_retries - 1:
                time.sleep(2 ** attempt)  # Exponential backoff: 1s, 2s, 4s
                continue
            
            logger.error(f"FATAL: Gemini AI Generation Error on attempt {attempt + 1}: {str(e)}")
            return jsonify({
                "success": False,
                "error": f"AI Service Error: {str(e)}"
            }), 503

# =====================================================================
# CHAT BOX
# =====================================================================
@app.route("/chat", methods=["POST"])
def chat():
    data = request.get_json()
    user_msg = data.get("message", "").strip().lower()
    if not user_msg:
        return jsonify({"reply": "Please type a message."})

    # Enhanced chat responses - Comprehensive knowledge base
    responses = {
        "hello": "Hello! Welcome to the ML Department portal. How can I help you today?",
        "hi": "Hi there! I'm here to assist you with information about our department.",
        "help": "I can assist with:\n• Department programs & details\n• Faculty information & expertise\n• Placement & job offers\n• Study topics (AI, ML, Data Science)\n• Resume templates (5 formats)\n• Pass percentage & attendance\n• Student resources\n• Previous year papers\n• Contact details\n\nWhat would you like to know?",
        "department": "We offer three specialized programs:\n• ACSML: Applied Computer Science and Machine Learning\n• NCSML: Neural Computing and Systems Machine Learning\n• DCSML: Data Science and Computational Machine Learning\n\nEach focuses on AI/ML with 100+ guided projects.",
        "faculty": "Our faculty specializes in:\n• Machine Learning & Deep Learning\n• Data Science & Analytics\n• AI & Computer Vision\n• Natural Language Processing\n• Cybersecurity & Cloud\n\nFaculty regularly publish in top conferences and guide research.",
        "faculty details": "For faculty profiles, visit Faculty section homepage. Each has research areas, publications, and office hours for consultations.",
        "student": "Students can:\n• Build professional resumes (5 templates)\n• Track CGPA, SGPA, attendance\n• View previous year papers\n• Access placement statistics\n• Manage project portfolios\n\nLogin to dashboard for full access.",
        "resume": "Our resume builder offers:\n• 5 professional templates\n• 100+ project descriptions\n• Career objective generator\n• Auto project enhancement\n• One-click PDF download\n\nChoose a template and start building!",
        "templates": "Resume Templates:\n📋 Classic ATS - Corporate/Finance\n✨ Modern Student - Startups\n💼 Minimal Tech - FAANG Tech\n🎓 Academic - Research positions\n🎨 Creative - Design-focused roles",
        "contact": "Reach us:\n• Email: info@ml-dept.edu\n• Phone: +91-9876543210\n• Address: Loyola Academy, Old Alwal, Secunderabad\n• Office Hours: Monday-Friday 9 AM - 5 PM",
        "admission": "Programs:\n• BTech CS (ML specialization)\n• MTech Machine Learning\n• Selection: Merit-based entrance exam\n• Visit our website for applications.",
        "placement": "Placement Stats:\n• 95%+ placement rate\n• Average: 8-12 LPA\n• Top recruiters: Google, Microsoft, Amazon, Goldman Sachs\n• Profiles: ML Engineer, Data Scientist, AI Developer\n• Internships: 4-6 months, 70%+ conversion to full-time",
        "pass percentage": "Department maintains 90%+ pass rate. Requirements:\n• 75%+ attendance for placement eligibility\n• 65%+ minimum for exam appearance\n• Below 65%: Shortage, special permission needed",
        "attendance": "Attendance Policy:\n• 75% minimum for placement eligibility\n• 65% minimum for exam appearance\n• Below 65%: Marked as shortage\n• Tracked for theory and practical courses",
        "job offers": "Offers from:\n• FAANG: Google, Amazon, Facebook, Apple, Netflix\n• Finance: Goldman Sachs, JP Morgan, Citibank\n• Consulting: McKinsey, BCG, EY\n• Startups: Unacademy, Byju's, OYO\n• Average: 8-12 LPA, 12-20 LPA for experienced",
        "internship": "Internships:\n• Duration: 4-6 months\n• Stipend: 30,000-60,000/month\n• Companies: Microsoft, Amazon, Samsung, Infosys\n• Profiles: ML, Data Analyst, Software Dev\n• 70%+ convert to full-time offers",
        "project": "Student projects:\n• ML algorithms & NLP\n• Data analysis & visualization\n• AI chatbots & systems\n• Image recognition & vision\n• Recommendation engines\n• Web apps & APIs\n• Research & publications",
        "machine learning": "Topics covered:\n• Supervised Learning (Regression, Classification)\n• Unsupervised Learning (Clustering, Reduction)\n• Deep Learning (CNNs, RNNs, Transformers)\n• NLP & Text Processing\n• Reinforcement Learning\n• Time Series Forecasting\n• Model Deployment",
        "ai": "AI Fundamentals:\n• Neural Networks & Deep Learning\n• Computer Vision (Image Classification, Detection)\n• NLP (Text Analysis, Chatbots)\n• Knowledge Representation\n• Search & Game Theory\n• Ethical AI & Responsible ML",
        "data science": "Data Science Topics:\n• Statistical Analysis & Testing\n• Data Cleaning & Preprocessing\n• Exploratory Data Analysis (EDA)\n• Predictive Modeling\n• Big Data (Spark, Hadoop)\n• BI & Dashboards\n• A/B Testing",
        "objective": "Career Objective Tips:\n• Highlight key strengths\n• Mention field of interest\n• Specify role type (Internship/Full-time)\n• Include 2-3 year goals\n\nExample: 'ML Engineer seeking to build scalable AI solutions and contribute to cutting-edge research projects.'",
        "career objective": "Strong objective includes:\n1. Your expertise\n2. Passion/interest area\n3. Goal (role/company type)\n4. Value you bring\n\nKeep concise, specific, 2-4 lines.",
        "resume objective": "Tailor to job! Include relevant skills from job description. Make it powerful, specific, under 4 lines.",
        "help objective": "Tell me: What's your main skill? What field? What type of role? I'll suggest a polished version!",
        "help project": "Format: '[Project Name] - [Tech] - [Impact]'\nExample: 'Movie Recommendation - Python, Filtering - 85% accuracy'\nOne per line, I'll enhance!",
        "help projects": "Enter projects with:\n• Project name\n• Technologies\n• Outcome/impact\n\nI'll auto-expand with professional descriptions!",
        "previous year papers": "Previous Year Papers:\n• Login to dashboard\n• Click 'Previous Year Papers' sidebar\n• Available for all semesters\n• All departments covered\n• Use for exam prep & practice",
        "papers": "Click 'Previous Year Papers' in sidebar to download all question papers for exam preparation!",
        "thank": "You're welcome! Ask anytime.",
        "bye": "Goodbye! Have a great day.",
        "goodbye": "Stay connected with our department!",
    }

    # Check for keywords in user message
    for keyword, response in responses.items():
        if keyword in user_msg:
            return jsonify({"reply": response})

    # Default responses for common queries
    if any(word in user_msg for word in ["what", "how", "when", "where", "why", "who"]):
        return jsonify({"reply": "I'd be happy to help with that information. Could you please be more specific about what you're looking for? You can ask about our programs, faculty, admissions, or any other department-related topics."})

    if any(word in user_msg for word in ["login", "password", "account"]):
        return jsonify({"reply": "For login issues or account help, please contact the department administrator or check the login page for instructions. Default password for demo accounts is '123456'."})

    # Generic helpful response
    return jsonify({"reply": "I'm here to help with information about the Machine Learning Department. You can ask me about our programs, faculty, student resources, or any other department-related questions. What would you like to know?"})

# This ensures the database is checked/created as soon as the app starts,
# regardless of whether it's run via 'python app.py' or a WSGI server like Gunicorn.

app = Flask(__name__)
