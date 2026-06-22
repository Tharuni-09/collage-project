from flask import Flask, flash, send_file, render_template, request, redirect, session, current_app, jsonify
from gemini_service import client
from google import genai
import os
import logging
import time
import requests
 
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass
 
# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)
 
 #-------------------------------------
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if GEMINI_API_KEY:
    client = genai.Client(api_key=GEMINI_API_KEY)
else:
    logger.warning("GEMINI_API_KEY not found. AI features will be disabled.")
    client = None
 
import sqlite3
import hashlib
import re
import io
from flask import url_for
from datetime import datetime, timedelta
from reportlab.lib.pagesizes import A4
from reportlab.platypus import (
    SimpleDocTemplate, BaseDocTemplate, PageTemplate, Frame,
    Paragraph, Spacer, Table, TableStyle, HRFlowable,
    FrameBreak, Image, KeepTogether, PageBreak
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch, mm
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY, TA_RIGHT
from werkzeug.utils import secure_filename
 
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH  = os.path.join(BASE_DIR, "database", "ml_dept.db")
 
# ── FIX 3: Create Flask app only ONCE (was created twice, second overwrote config) ──
app = Flask(
    __name__,
    static_folder=os.path.join(os.path.dirname(__file__), "static")
)
app.secret_key = os.environ.get("SECRET_KEY", "change_me_in_production")
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(days=365)
app.config['SESSION_COOKIE_PERMANENT']   = True
app.config['SESSION_COOKIE_SECURE']      = False   # Set True in production (HTTPS)
app.config['SESSION_COOKIE_HTTPONLY']    = True
app.config['SESSION_COOKIE_SAMESITE']   = 'Lax'
app.config['UPLOAD_FOLDER']             = os.path.join(BASE_DIR, 'static', 'notes')
app.config['MAX_CONTENT_LENGTH']        = 50 * 1024 * 1024  # 50 MB
 
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
 
# ── FIX 4: Remove bare sqlite3.connect() call that ran at module level ───────
#   (was: conn = sqlite3.connect("collage.db"); cursor = ...; print(...))
 
# =====================================================================
# HELPER FUNCTIONS
# =====================================================================
 
def should_print_section(text, user_response):
    if not text or not text.strip():
        return False
    low = user_response.strip().lower()
    if low in ["no", "false", "0", "none", "not applicable", "na"]:
        return False
    return True
 
 
def elaborate(text, section_name):
    if not text or section_name not in ["skills", "projects", "internships"]:
        return text
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    if not lines:
        return ""
    if section_name == "skills":
        bullets = []
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.lower().startswith("python"):
                bullets.append("Python: Proficient in core Python, OOP, and popular libraries such as NumPy and Pandas.")
            elif line.lower().startswith("flask"):
                bullets.append("Flask: Developed web applications and REST APIs using Flask and SQLite.")
            elif line.lower().startswith("machine learning") or "ml" in line.lower():
                bullets.append("Machine Learning: Experience with classification, regression, and neural networks using scikit-learn or TensorFlow/PyTorch.")
            elif "sql" in line.lower() or "database" in line.lower():
                bullets.append("SQL & Databases: Designed and queried SQLite and MySQL databases for web applications.")
            elif "html" in line.lower() and "css" in line.lower():
                bullets.append("Web Technologies: Built responsive websites using HTML, CSS, and JavaScript (or frameworks).")
            else:
                bullets.append(f"{line}: Comfortable applying this technology to real-world projects.")
        return "\n".join(bullets)
    if section_name == "projects":
        bullets = []
        for i, line in enumerate(lines):
            if i == 0 and len(lines) == 1:
                bullets.append(f"{line} – A well-structured project applying core concepts of the technology stack.")
            elif i == 0:
                bullets.append(f"Project: {line}")
            else:
                bullets.append(f"• {line}")
        return "\n".join(bullets)
    if section_name == "internships":
        bullets = []
        for line in lines:
            if line.strip().lower().startswith("intern") or "internship" in line.lower():
                bullets.append("Internship role focused on practical software development in a team environment.")
            else:
                bullets.append(f"• {line}")
        return "\n".join(bullets)
    return text
 
 
def polish_objective(obj):
    if not obj.strip():
        return ""
    words = obj.strip().split()
    if not words:
        return ""
    obj = " ".join(words)
    if not obj.endswith("."):
        obj += "."
    return f"Motivated and detail-oriented professional with strong {obj[0].lower() + obj[1:]}"
 
 
def expand_section(text, section_name):
    if not text or text.lower().strip() in ["no", "none", "na", "false", "0"]:
        return ""
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if not lines:
        return ""
    expanded = []
    if section_name == "objective":
        return (
            "Accomplished professional with a background in "
            f"{text}. Dedicated to leveraging technical skills to drive operational "
            "excellence and contribute to organizational success in a challenging role."
        )
    elif section_name == "work_experience":
        for line in lines:
            expanded.append(f"• {line}")
            expanded.append("  - Achieved key deliverables by optimizing workflows and improving efficiency by 15–20%.")
            expanded.append("  - Collaborated with cross-functional teams to resolve complex technical challenges.")
        return "\n".join(expanded)
    elif section_name == "projects":
        for line in lines:
            expanded.append(f"• Project: {line}")
            expanded.append("  - Developed core architecture using modern frameworks, ensuring scalability and high performance.")
            expanded.append("  - Conducted thorough testing and debugging, resulting in a 10% improvement in system stability.")
        return "\n".join(expanded)
    return "\n".join([f"• {l}" for l in lines])
 
 
COMMON_SPELLING_FIXES = {
    "oblective": "objective", "doutd": "doubt", "recieve": "receive",
    "teh": "the", "adn": "and", "acheivement": "achievement",
    "deatils": "details", "develpment": "development", "manger": "manager",
    "projct": "project", "programe": "program", "speling": "spelling",
    "currect": "correct",
}
 
 
def clean_text(text):
    if not text:
        return ""
    text = str(text).strip()
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    text = text.replace("\n", "<br/>")
    for wrong, right in COMMON_SPELLING_FIXES.items():
        text = re.sub(rf"\b{wrong}\b", right, text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text)
    return text
 
 
def split_items(text):
    if not text:
        return []
    return [i.strip() for i in re.split(r'[\n,;]+', str(text)) if i.strip()]
 
 
def enhance_objective(objective, field_of_study):
    objective = clean_text(objective)
    if not objective:
        field = field_of_study or "Machine Learning"
        return (f"Results-oriented {field} professional with a focus on delivering high-impact technical "
                "solutions through advanced machine learning methodologies and scalable software architecture.")
    if len(objective.split()) < 15:
        objective = objective.rstrip('.')
        objective += ", eager to contribute to innovative teams and dedicated to leveraging data-driven insights to solve complex challenges."
    return objective
 
 
def parse_projects(projects_text):
    if not projects_text or not projects_text.strip():
        return []
    return [clean_text(line) for line in projects_text.splitlines() if line.strip()]
 
 
def describe_project(project_line):
    return []
 
 
def describe_experience(exp_line):
    text  = exp_line.strip()
    lower = text.lower()
    bullets = []
    if any(w in lower for w in ["software", "developer", "sde", "engineer", "backend", "frontend"]):
        bullets.append(f"Contributed to high-quality codebases as a {text}, focusing on scalable features and system reliability.")
        bullets.append("Collaborated with cross-functional teams to deliver robust software solutions using modern frameworks and best practices.")
        bullets.append("Optimized application performance and resolved complex technical bottlenecks in production environments.")
    elif any(w in lower for w in ["data", "analyst", "scientist", "ml", "machine learning"]):
        bullets.append(f"Leveraged statistical methodologies in my role as {text} to extract actionable insights from complex datasets.")
        bullets.append("Developed automated reporting pipelines and visualization dashboards to drive data-driven decision making.")
        bullets.append("Performed deep-dive exploratory data analysis to identify trends and improve model predictive accuracy.")
    elif any(w in lower for w in ["research", "intern", "academic"]):
        bullets.append(f"Conducted in-depth research for {text}, contributing to critical project milestones and technical documentation.")
        bullets.append("Documented technical findings and presented comprehensive reports to key stakeholders and mentors.")
        bullets.append("Assisted in the design and execution of experimental workflows to validate theoretical hypotheses.")
    else:
        bullets.append(f"Supported organizational goals as {text} through effective task management and cross-team communication.")
        bullets.append("Developed professional domain expertise while contributing to organizational efficiency and project success.")
        bullets.append("Maintained high standards of quality and reliability in all assigned deliverables.")
    return bullets
 
 
def describe_certification(cert_line):
    text  = cert_line.strip()
    if not text:
        return ""
    lower = text.lower()
    if any(w in lower for w in ["aws", "google", "azure", "cloud"]):
        return "Validated cloud architecture expertise and proficiency in designing secure, scalable distributed systems."
    if any(w in lower for w in ["python", "java", "sql", "javascript"]):
        return f"Demonstrated advanced technical proficiency and problem-solving capabilities in {text}."
    return f"Professional credential validating advanced proficiency and theoretical knowledge in {text}."
 
 
# =====================================================================
# DATABASE
# =====================================================================
 
# ── FIX 5: Single get_db() function using the correct ML-dept DB path ────────
def get_db():
    os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
    conn = sqlite3.connect(DB_PATH, timeout=10)
    conn.row_factory = sqlite3.Row
    return conn
 
 
# =====================================================================
# INIT DB SCHEMA
# =====================================================================
 
def init_database():
    db = get_db()
    try:
        table_exists = db.execute(
            "SELECT name FROM sqlite_master WHERE type='table' AND name='users'"
        ).fetchone()
        if not table_exists:
            with app.open_resource("schema.sql") as f:
                db.executescript(f.read().decode("utf-8"))
 
        try:
            db.execute("ALTER TABLE users ADD COLUMN is_hod INTEGER DEFAULT 0")
        except Exception:
            pass
 
        db.execute("""
            CREATE TABLE IF NOT EXISTS previous_year_papers (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                subject TEXT NOT NULL,
                branch TEXT NOT NULL,
                year INTEGER NOT NULL,
                semester INTEGER NOT NULL,
                photo_limit INTEGER DEFAULT 1
            )
        """)
        db.execute("""
            CREATE TABLE IF NOT EXISTS paper_photos (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                paper_id INTEGER NOT NULL,
                image_path TEXT NOT NULL,
                FOREIGN KEY (paper_id) REFERENCES previous_year_papers(id)
            )
        """)
        # Notes tables
        db.execute("""
            CREATE TABLE IF NOT EXISTS departments (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                name TEXT UNIQUE NOT NULL
            )
        """)

        db.execute("""
            CREATE TABLE IF NOT EXISTS faculty_permissions (
                faculty_id INTEGER PRIMARY KEY,
                upload_notes INTEGER DEFAULT 1,
                upload_papers INTEGER DEFAULT 1,
                manage_students INTEGER DEFAULT 0,
                use_ai INTEGER DEFAULT 0
            )
        """)
        for dept in ["ACSML", "NCSML", "DCSML"]:
            try:
                db.execute("INSERT OR IGNORE INTO departments (name) VALUES (?)", (dept,))
            except Exception:
                pass

        db.execute("""
            CREATE TABLE IF NOT EXISTS notes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                title TEXT NOT NULL,
                description TEXT,
                file_path TEXT NOT NULL,
                file_name TEXT NOT NULL,
                department_id INTEGER NOT NULL,
                faculty_id TEXT NOT NULL,
                uploaded_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (department_id) REFERENCES departments(id)
            )
        """)

        db.execute("""
            CREATE TABLE IF NOT EXISTS admin_settings(
                id INTEGER PRIMARY KEY,
                resume_builder INTEGER DEFAULT 1,
                ppt_generator INTEGER DEFAULT 1,
                student_signup INTEGER DEFAULT 1,
                faculty_upload INTEGER DEFAULT 1,
                ai_resume_builder INTEGER DEFAULT 1
            )
            """)

        db.execute("""
            INSERT OR IGNORE INTO admin_settings(
                id,
                resume_builder,
                ppt_generator,
                student_signup,
                faculty_upload,
                ai_resume_builder
            )
            VALUES(1,1,1,1,1,1)
            """)
        db.commit()
    finally:
        db.close()
 
# =====================================================================
# INDEX AND STATIC PAGES
# =====================================================================
 
@app.route("/")
def index():
    latest_outreach = None
    try:
        conn = get_db()
        row = conn.execute(
            "SELECT id, title, year, venue, participants, description, photo "
            "FROM outreach_programs ORDER BY id DESC LIMIT 1"
        ).fetchone()
        conn.close()
        if row:
            latest_outreach = {
                "id": row[0], "title": row[1], "year": row[2],
                "venue": row[3], "participants": row[4],
                "description": row[5], "photo": row[6]
            }
    except Exception as e:
        logger.error(f"Error fetching latest outreach: {e}")
    return render_template("index.html", latest_outreach=latest_outreach)
 
 
@app.route("/outreach")
def outreach():
    return render_template("outreach.html")
 
 
@app.route("/gallery")
def gallery():
    return render_template("gallery.html")
 
 
# ── FIX 6: Added /chat route (was referenced in base.html JS but missing) ────
@app.route("/chat", methods=["POST"])
def chat():
    data    = request.get_json() or {}
    message = data.get("message", "").strip()
    if not message:
        return jsonify({"reply": "Please type a message."})
        
    if not client:
        return jsonify({"reply": "⚠️ **Configuration Error**: The AI is currently offline. Please ensure the `GEMINI_API_KEY` environment variable is set in the Render deployment dashboard."})
        
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=(
                "You are an advanced, knowledgeable AI assistant. "
                "You can answer ANY question on ANY topic — academics, science, technology, "
                "programming, mathematics, history, current affairs, career advice, general knowledge, "
                "creative writing, code debugging, and more. "
                "You are also the official assistant for the ML Department Portal of Loyola Academy, "
                "so you know about courses, notes, resumes, outreach programs, and department activities. "
                "Provide clear, accurate, well-structured answers. Use bullet points or numbered lists "
                "when helpful. Keep responses concise but thorough. "
                "If you are unsure, say so honestly.\n\nUser: " + message
            )
        )
        return jsonify({"reply": response.text})
    except Exception as e:
        logger.error(f"Chat Gemini error: {e}")
        return jsonify({"reply": "Sorry, I'm having trouble responding right now. Please try again later."})


 #----------------------------------------------------------
@app.route("/admin/change-password", methods=["GET", "POST"])
def admin_change_password():

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    if request.method == "POST":

        username = request.form.get("username")
        password = request.form.get("password")

        import hashlib

        password_hash = hashlib.sha256(
            password.encode()
        ).hexdigest()

        db = get_db()

        db.execute("""
        UPDATE users
        SET username=?,
            password_hash=?
        WHERE role='admin'
        """,
        (
            username,
            password_hash
        ))

        db.commit()
        db.close()

        return render_template(
            "admin/change_password.html",
            success="Admin Credentials Updated Successfully"
        )

    return render_template("admin/change_password.html")
# -------------------OUTREACH--------------------------------------
 
@app.route('/add_outreach_program', methods=['POST'])
def add_outreach_program():
    if session.get('role') != 'faculty':
        return "Unauthorized", 403
    title        = request.form['title']
    year         = request.form['year']
    venue        = request.form['venue']
    participants = request.form['participants']
    description  = request.form['description']
    file         = request.files['photo']
    filename     = secure_filename(file.filename)
    file.save(os.path.join('static/outreach', filename))
    conn = get_db()
    conn.execute("""
        INSERT INTO outreach_programs (title, year, venue, participants, description, photo)
        VALUES (?, ?, ?, ?, ?, ?)
    """, (title, year, venue, participants, description, filename))
    conn.commit()
    conn.close()
    return redirect(url_for('outreach'))
 
 
# -------------------Previous Year Papers-----------------------
 
def init_paper_grid_db():
    conn = get_db()
    conn.execute("""
        CREATE TABLE IF NOT EXISTS paper_grid (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            regulation TEXT NOT NULL,
            year TEXT NOT NULL,
            semester TEXT NOT NULL,
            subject TEXT NOT NULL
        )
    """)
    conn.commit()
    conn.close()
 
 
@app.route("/add_paper_grid", methods=["POST"])
def add_paper_grid():
    regulation = request.form.get("regulation")
    year       = request.form.get("year")
    semester   = request.form.get("semester")
    subject    = request.form.get("subject")
    conn = get_db()
    conn.execute("""
        INSERT INTO paper_grid (regulation, year, semester, subject)
        VALUES (?, ?, ?, ?)
    """, (regulation, year, semester, subject))
    conn.commit()
    conn.close()
    return redirect(url_for("previous_year_papers"))
 
 
@app.route("/previous-year-papers")
def previous_year_papers():
    db     = get_db()
    papers = []
    # ── FIX 7: was querying "papers" table (wrong name); correct table is "previous_year_papers"
    papers_rows = db.execute("SELECT * FROM previous_year_papers").fetchall()
    for row in papers_rows:
        paper  = dict(row)
        photos = db.execute(
            "SELECT image_path FROM paper_photos WHERE paper_id = ?", [paper["id"]]
        ).fetchall()
        paper["photos"] = [p["image_path"] for p in photos]
        papers.append(paper)
    db.close()
    return render_template("previous_year_papers.html", session=session, papers=papers)
 
 
@app.route("/upload-paper-camera", methods=["POST"])
def upload_paper_camera():
    if not session.get("uid"):
        return jsonify({"success": False, "message": "Login required"})
    paper_id   = request.form.get("paper_id")
    image_data = request.files.get("image")
    if image_data and paper_id:
        db    = get_db()
        paper = db.execute(
            "SELECT photo_limit FROM previous_year_papers WHERE id = ?", [paper_id]
        ).fetchone()
        current_count = db.execute(
            "SELECT COUNT(*) as count FROM paper_photos WHERE paper_id = ?", [paper_id]
        ).fetchone()
        if paper and current_count["count"] >= paper["photo_limit"]:
            db.close()
            return jsonify({"success": False,
                            "message": f"Maximum photo limit ({paper['photo_limit']}) reached."})
        filename = f"paper_{paper_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.jpg"
        filepath = os.path.join(app.static_folder, "images", filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        image_data.save(filepath)
        db.execute("INSERT INTO paper_photos (paper_id, image_path) VALUES (?, ?)", [paper_id, filename])
        db.commit()
        db.close()
        return jsonify({"success": True, "filename": filename})
    return jsonify({"success": False})
 
 
# =====================================================================
# LOGIN
# =====================================================================
 
@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "GET":
        message = request.args.get("message")
        return render_template("login.html", error=None, message=message)
 
    uid      = request.form.get("uid") or request.form.get("username") or request.form.get("identifier")
    password = request.form.get("password")
    role     = request.form.get("role") or request.form.get("mode")
 
    is_ajax = request.is_json or request.headers.get("X-Requested-With") == "XMLHttpRequest"
    if is_ajax and request.is_json:
        data     = request.get_json() or {}
        uid      = data.get("uid") or data.get("username") or uid
        password = data.get("password") or password
        role     = data.get("role") or data.get("mode") or role
 
    uid     = str(uid).strip() if uid else None
    role    = str(role).strip().lower() if role else None
    uid_int = int(uid) if uid and uid.isdigit() else None
 
    db  = get_db()
    row = db.execute("""
        SELECT u.uid, u.username, u.password_hash, u.role, u.name,
               u.email, u.phone, u.is_hod, s.roll
        FROM users u
        LEFT JOIN students s ON u.uid = s.uid
        WHERE u.uid = ? OR u.uid = ? OR u.username = ? OR u.email = ? OR s.roll = ?
    """, [uid, uid_int, uid, uid, uid]).fetchone()
 
    if not row:
        db.close()
        msg = "Invalid UID or ID not found. Please check your credentials."
        if is_ajax:
            return jsonify({"success": False, "message": msg})
        return render_template("login.html", error=msg)
 
    db_role      = row["role"]
    user_is_hod  = int(row["is_hod"] or 0) == 1
    db_uid       = row["uid"]
    db_hash      = row["password_hash"]
    full_name    = row["name"]
 
    entered_hash = hashlib.sha256((password or "").encode()).hexdigest()
    if db_hash != entered_hash:
        db.close()
        if is_ajax:
            return jsonify({"success": False, "message": "Invalid password"})
        return render_template("login.html", error="Invalid password")
 
    if role == "student" and db_role != "student":
        db.close()
        msg = "This ID is registered as Faculty/HOD. Please select the correct role."
        if is_ajax:
            return jsonify({"success": False, "message": msg})
        return render_template("login.html", error=msg)
 
    if role == "faculty" and db_role != "faculty":
        db.close()
        msg = "This ID is registered as a Student. Please select the correct role."
        if is_ajax:
            return jsonify({"success": False, "message": msg})
        return render_template("login.html", error=msg)
        
    if role == "admin" and db_role != "admin":
        db.close()
        msg = "This ID does not have Admin privileges."
        if is_ajax:
            return jsonify({"success": False, "message": msg})
        return render_template("login.html", error=msg)
 
    db.close()
    session["uid"]    = db_uid
    session["role"]   = db_role
    session["name"]   = full_name
    session["is_hod"] = user_is_hod
    session.permanent = True
 
    if is_ajax:
        return jsonify({"success": True, "message": "Login successful!", "role": db_role})
    if db_role == "student":
        return redirect(url_for("student_dashboard"))

    elif db_role == "admin":
        return redirect(url_for("admin_dashboard"))
    return redirect(url_for("faculty_dashboard"))
 
 #admin antaru ra babu
@app.route("/admin-login", methods=["GET", "POST"])
def admin_login():

    if request.method == "GET":
        return render_template("admin/admin_login.html")

    username = request.form.get("username")
    password = request.form.get("password")

    db = get_db()

    admin = db.execute("""
        SELECT *
        FROM users
        WHERE username = ?
        AND role = 'admin'
    """, [username]).fetchone()

    db.close()

    if not admin:
        return render_template(
            "admin/admin_login.html",
            error="Invalid Admin Credentials"
        )

    entered_hash = hashlib.sha256(
        password.encode()
    ).hexdigest()

    if entered_hash != admin["password_hash"]:
        return render_template(
            "admin/admin_login.html",
            error="Invalid Admin Credentials"
        )

    session["uid"] = admin["uid"]
    session["role"] = "admin"
    session["name"] = admin["name"]

    return redirect(url_for("admin_dashboard"))

@app.route("/admin-dashboard")
def admin_dashboard():

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    db = get_db()

    total_students = db.execute(
        "SELECT COUNT(*) FROM students"
    ).fetchone()[0]

    total_faculty = db.execute(
        "SELECT COUNT(*) FROM users WHERE role='faculty'"
    ).fetchone()[0]

    total_resumes = db.execute(
        "SELECT COUNT(*) FROM resumes"
    ).fetchone()[0]

    total_notes = db.execute(
        "SELECT COUNT(*) FROM notes"
    ).fetchone()[0]

    db.close()

    return render_template(
        "admin/admin_dashboard.html",
        total_students=total_students,
        total_faculty=total_faculty,
        total_resumes=total_resumes,
        total_notes=total_notes
    )

@app.route("/admin/students")
def admin_students():

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    db = get_db()

    students = db.execute("""
        SELECT uid, name, username
        FROM users
        WHERE role='student'
    """).fetchall()

    db.close()

    return render_template(
        "admin/manage_students.html",
        students=students
    )

@app.route("/admin/faculty")
def admin_faculty():
    
    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    db = get_db()

    faculty = db.execute("""
        SELECT
            uid,
            name,
            username
        FROM users
        WHERE role='faculty'
    """).fetchall()

    db.close()

    return render_template(
        "admin/manage_faculty.html",
        faculty=faculty
    )

@app.route("/admin/resumes")
def admin_resumes():

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    db = get_db()

    resumes = db.execute("""
        SELECT
            id,
            uid,
            student_id,
            title,
            created_at
        FROM resumes
        ORDER BY created_at DESC
    """).fetchall()

    db.close()

    return render_template(
        "admin/manage_resumes.html",
        resumes=resumes
    )

@app.route("/admin/notes")
def admin_notes():

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    db = get_db()

    notes = db.execute("""
        SELECT
            id,
            title,
            file_name,
            uploaded_at
        FROM notes
        ORDER BY uploaded_at DESC
    """).fetchall()

    db.close()

    return render_template(
        "admin/manage_notes.html",
        notes=notes
    )

@app.route("/admin/papers")
def admin_papers():

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    db = get_db()

    papers = db.execute("""
        SELECT
            id,
            subject,
            year,
            semester,
            image_path
        FROM previous_year_papers
        ORDER BY year DESC
    """).fetchall()

    db.close()

    return render_template(
        "admin/manage_papers.html",
        papers=papers
    )

@app.route("/check-papers")
def check_papers():

    db = get_db()

    papers = db.execute("""
        SELECT *
        FROM previous_year_papers
    """).fetchall()

    db.close()

    return str([dict(x) for x in papers])

@app.route("/admin/settings")
def admin_settings():

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    db = get_db()

    settings = db.execute("""
    SELECT *
    FROM admin_settings
    WHERE id = 1
    """).fetchone()

    db.close()

    return render_template(
        "admin/settings.html",
        settings=settings
    )

@app.route("/admin/toggle-setting/<setting>")
def toggle_setting(setting):

    if session.get("role") != "admin":
        return redirect(url_for("admin_login"))

    allowed = [
        "resume_builder",
        "ppt_generator",
        "student_signup",
        "faculty_upload",
        "ai_resume_builder"
    ]

    if setting not in allowed:
        return "Invalid Setting"

    db = get_db()

    current = db.execute(
        f"SELECT {setting} FROM admin_settings WHERE id=1"
    ).fetchone()[0]

    new_value = 0 if current else 1

    db.execute(
        f"UPDATE admin_settings SET {setting}=? WHERE id=1",
        (new_value,)
    )

    db.commit()
    db.close()

    return redirect("/admin/settings")

# =====================================================================
# FORGOT PASSWORD
# =====================================================================
 
@app.route("/forgot-password", methods=["GET", "POST"])
def forgot_password():
    if request.method == "GET":
        return render_template("forget_password.html", error=None, step="request")
 
    action = request.form.get("action", "request")
    if action == "request":
        identifier = clean_text(request.form.get("identifier", ""))
        if not identifier:
            return render_template("forget_password.html", error="Missing identifier", step="request")
        db  = get_db()
        row = db.execute("""
            SELECT u.uid, u.email, u.phone, s.roll
            FROM users u
            LEFT JOIN students s ON u.uid = s.uid
            WHERE u.uid = ? OR u.email = ? OR u.phone = ? OR s.roll = ?
        """, [identifier, identifier, identifier, identifier]).fetchone()
        db.close()
        if not row:
            return render_template("forget_password.html",
                                   error="No user found with that UID, email, phone, or roll number.",
                                   step="request")
        user_email = row["email"] if isinstance(row, sqlite3.Row) else row[1]
        if not user_email:
            return render_template("forget_password.html",
                                   error="No email address is available for this account.",
                                   step="request")
        return render_template("forget_password.html", step="reset", uid=row["uid"],
                               email=user_email,
                               message=f"A reset email would be sent to {user_email}. Enter your new password below.")
 
    if action == "reset":
        uid              = request.form.get("uid")
        email            = clean_text(request.form.get("email", ""))
        password         = request.form.get("password", "")
        password_confirm = request.form.get("password_confirm", "")
        if not all([uid, email, password, password_confirm]):
            return render_template("forget_password.html", error="All fields are required.",
                                   step="reset", uid=uid, email=email)
        if password != password_confirm:
            return render_template("forget_password.html", error="Passwords do not match.",
                                   step="reset", uid=uid, email=email)
        db  = get_db()
        row = db.execute("SELECT email FROM users WHERE uid = ?", [uid]).fetchone()
        if not row or (row["email"] if isinstance(row, sqlite3.Row) else row[0]) != email:
            db.close()
            return render_template("forget_password.html", error="Email verification failed.", step="request")
        password_hash = hashlib.sha256(password.encode()).hexdigest()
        db.execute("UPDATE users SET password_hash = ? WHERE uid = ?", [password_hash, uid])
        db.commit()
        db.close()
        return render_template("forget_password.html",
                               message="Password reset successfully. You can now login with your new password.",
                               step="done")
 
 
# =====================================================================
# STUDENT DASHBOARD
# =====================================================================
 
@app.route("/student-dashboard")
def student_dashboard():
    uid = session.get("uid")
    if not uid:
        return redirect(url_for("login"))
    if session.get("role") != "student":
        return redirect(url_for("login"))
 
    db  = get_db()
    row = db.execute("""
        SELECT s.roll, s.name, s.department, s.cgpa, s.sgpa, s.attendance
        FROM students s
        WHERE s.uid = ?
    """, [uid]).fetchone()
 
    if not row:
        db.close()
        return render_template("student_dashboard.html", roll="N/A", name="Unknown Student",
                               department="N/A", cgpa=0, sgpa=0, attendance=0, resumes=[],
                               error="Student profile not found. Please contact the administrator."), 404
 
    roll, name, department, cgpa, sgpa, attendance = row
 
    resumes = db.execute("""
        SELECT r.id, r.uid, r.title, r.created_at
        FROM resumes r
        WHERE r.uid = ?
        ORDER BY r.created_at DESC
    """, [uid]).fetchall()
    db.close()
 
    return render_template("student_dashboard.html", roll=roll, name=name, department=department,
                           cgpa=cgpa, sgpa=sgpa, attendance=attendance, resumes=resumes, session=session)
 
 
@app.route("/resume_form")
def resume_form():
    if "uid" not in session:
        return redirect(url_for("login"))
    return render_template("resume_form.html", session=session)
 
 
@app.route("/generate", methods=["POST"])
def generate():
    form_data = request.form.to_dict(flat=False)
    for key, value in form_data.items():
        if len(value) == 1:
            form_data[key] = value[0]
 
    template_choice = form_data.get("template", "classic")
 
    photo = request.files.get("profile_photo")
    if photo and photo.filename:
        out_dir        = os.path.join(BASE_DIR, "output")
        os.makedirs(out_dir, exist_ok=True)
        uid_prefix     = session.get('uid', 'anon')
        photo_filename = f"temp_{uid_prefix}_{photo.filename}"
        photo_path     = os.path.join(out_dir, photo_filename)
        photo.save(photo_path)
        form_data['photo_path'] = photo_path
 
    pdf_content = generate_resume_pdf(template_choice, form_data)
 
    uid = session.get("uid")
    if uid:
        db = get_db()
        existing_resume = db.execute(
            "SELECT id FROM resumes WHERE uid = ?",
            [uid]
        ).fetchone()
        
        new_title = f"Resume ({template_choice.capitalize()})"
        if existing_resume:
            db.execute("UPDATE resumes SET title = ?, pdf_content = ?, created_at = ? WHERE id = ?",
                       [new_title, pdf_content, datetime.now().strftime('%Y-%m-%d %H:%M:%S'), existing_resume["id"]])
        else:
            db.execute(
                "INSERT INTO resumes (uid, student_id, title, pdf_content, created_at) VALUES (?, ?, ?, ?, ?)",
                [uid, uid, new_title, pdf_content, datetime.now().strftime('%Y-%m-%d %H:%M:%S')]
            )
        db.commit()
        db.close()
 
    filename = f"{clean_text(form_data.get('name') or 'resume').replace(' ', '_')}_{template_choice}.pdf"
 
    if 'photo_path' in form_data and os.path.exists(form_data['photo_path']):
        try:
            os.remove(form_data['photo_path'])
        except Exception:
            pass
 
    return send_file(io.BytesIO(pdf_content), as_attachment=False,
                     download_name=filename, mimetype="application/pdf")
 
 
# ── FIX 8: Added /ppt_generator route (was missing entirely) ─────────────────
@app.route("/ppt_generator")
def ppt_generator():
    if "uid" not in session:
        return redirect(url_for("login"))
    return render_template("ppt_generator.html", session=session)
 
 
# ── FIX 9: Added /generate-ppt route (called from ppt_generator.html) ────────
@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    if "uid" not in session:
        return jsonify({"success": False, "error": "Login required"}), 401
 
    file   = request.files.get("file")
    slides = request.form.get("slides", 20)
 
    if not file:
        return jsonify({"success": False, "error": "No file uploaded"}), 400
 
    # Read file content for the prompt
    file_content = ""
    try:
        filename_lower = file.filename.lower()
        if filename_lower.endswith(".txt"):
            file_content = file.read().decode("utf-8", errors="ignore")
        elif filename_lower.endswith(".pdf"):
            import pdfplumber
            with pdfplumber.open(file) as pdf:
                file_content = "\n".join(page.extract_text() or "" for page in pdf.pages[:20])
        else:
            file_content = file.read().decode("utf-8", errors="ignore")
    except Exception as e:
        logger.warning(f"File reading error: {e}")
        file_content = "Unable to extract text from file."
 
    prompt = f"""
Create a structured outline for a {slides}-slide PowerPoint presentation based on the following content.
 
Content:
{file_content[:6000]}
 
Output format (JSON array):
[
  {{"slide": 1, "title": "...", "points": ["...", "...", "..."]}},
  ...
]
 
Rules:
- Return ONLY valid JSON, no markdown
- Max 5 bullet points per slide
- Keep points concise (under 15 words each)
- Include title slide, agenda, content slides, and conclusion
"""
    try:
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        raw      = response.text.strip().lstrip("```json").lstrip("```").rstrip("```").strip()
 
        # Build a real PPTX using python-pptx
        import json
        from pptx import Presentation
        from pptx.util import Inches, Pt
 
        slides_data = json.loads(raw)
        session["ppt_review_data"] = slides_data
        return jsonify({
            "success": True,
            "review_url": url_for("ppt_review")
            })
 
        for slide_info in slides_data:
            layout     = prs.slide_layouts[1]
            pptx_slide = prs.slides.add_slide(layout)
            title_ph   = pptx_slide.shapes.title
            body_ph    = pptx_slide.placeholders[1]
 
            title_ph.text = slide_info.get("title", "")
            tf = body_ph.text_frame
            tf.clear()
            for point in slide_info.get("points", []):
                p = tf.add_paragraph()
                p.text  = point
                p.level = 0
 
        pptx_buf = io.BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)
 
        # Save to static for download
        out_filename = f"ppt_{session.get('uid', 'anon')}_{int(time.time())}.pptx"
        out_path     = os.path.join(BASE_DIR, "static", "uploads", "generated", out_filename)
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        with open(out_path, "wb") as f:
            f.write(pptx_buf.getvalue())
 
        download_url = url_for("static", filename=f"uploads/generated/{out_filename}")
        return jsonify({"success": True, "slides_data": slides_data})
 
    except Exception as e:
        logger.error(f"PPT generation error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500
    
@app.route("/ppt-review")
def ppt_review():

    if "ppt_review_data" not in session:
        return redirect("/faculty")

    slides_data = session["ppt_review_data"]

    return render_template(
        "ppt_review.html",
        slides=slides_data
    )
 
 
# ── FIX 10: Added /generate-notes route (called from ppt_generator.html) ─────
@app.route("/generate-notes", methods=["POST"])
def generate_notes_route():
    if "uid" not in session:
        return jsonify({"success": False, "error": "Login required"}), 401
 
    file     = request.files.get("file")
    sections = request.form.get("sections", 10)
 
    if not file:
        return jsonify({"success": False, "error": "No file uploaded"}), 400
 
    file_content = ""
    try:
        filename_lower = file.filename.lower()
        if filename_lower.endswith(".txt"):
            file_content = file.read().decode("utf-8", errors="ignore")
        else:
            file_content = file.read().decode("utf-8", errors="ignore")
    except Exception as e:
        logger.warning(f"File reading error: {e}")
        file_content = "Unable to extract text."
 
    prompt = f"""
Create comprehensive study notes with {sections} sections based on the following content.
 
Content:
{file_content[:6000]}
 
Output format (JSON):
{{
  "title": "...",
  "sections": [
    {{"heading": "...", "content": "...", "key_points": ["...", "..."]}},
    ...
  ]
}}
 
Rules:
- Return ONLY valid JSON, no markdown
- Each section should have a clear heading, paragraph content, and 3-5 key points
- Keep it academic and easy to understand
"""
    try:
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        raw      = response.text.strip().lstrip("```json").lstrip("```").rstrip("```").strip()
 
        import json
        notes_data = json.loads(raw)
 
        # Build a text/HTML notes file
        html_content = f"""<!DOCTYPE html>
<html lang='en'>
<head><meta charset='UTF-8'><title>{notes_data.get('title','Notes')}</title>
<style>
  body {{ font-family: Arial, sans-serif; max-width: 900px; margin: 40px auto; padding: 0 20px; color: #222; }}
  h1 {{ color: #0a3d62; border-bottom: 2px solid #0a3d62; }}
  h2 {{ color: #1e6fa8; margin-top: 30px; }}
  .key-points {{ background: #f0f7ff; border-left: 4px solid #1e6fa8; padding: 10px 16px; margin-top: 10px; }}
  .key-points li {{ margin: 4px 0; }}
</style></head>
<body>
<h1>{notes_data.get('title','Study Notes')}</h1>
"""
        for sec in notes_data.get("sections", []):
            html_content += f"<h2>{sec.get('heading','')}</h2>\n"
            html_content += f"<p>{sec.get('content','')}</p>\n"
            kp = sec.get("key_points", [])
            if kp:
                html_content += "<div class='key-points'><strong>Key Points:</strong><ul>"
                for point in kp:
                    html_content += f"<li>{point}</li>"
                html_content += "</ul></div>\n"
        html_content += "</body></html>"
 
        out_filename = f"notes_{session.get('uid', 'anon')}_{int(time.time())}.html"
        out_path     = os.path.join(BASE_DIR, "static", "uploads", "generated", out_filename)
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(html_content)
 
        download_url = url_for("static", filename=f"uploads/generated/{out_filename}")
        return jsonify({"success": True, "notes_data": notes_data})
 
    except Exception as e:
        logger.error(f"Notes generation error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500
 
 
# =====================================================================
# FACULTY ROUTES
# =====================================================================
 
@app.route("/faculty/<dept>")
def faculty_department(dept):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    db       = get_db()
    students = db.execute("""
        SELECT s.*, u.email, u.phone
        FROM students s
        JOIN users u ON s.uid = u.uid
        WHERE s.department = ?
        ORDER BY s.roll
    """, [dept]).fetchall()
    resumes  = db.execute("""
        SELECT r.id, r.uid, r.title, r.created_at, s.roll, s.name
        FROM resumes r
        JOIN students s ON r.uid = s.uid
        WHERE s.department = ?
        ORDER BY s.roll
    """, [dept]).fetchall()
    from collections import defaultdict
    grouped = defaultdict(list)
    for r in resumes:
        grouped[r['uid']].append(dict(r))
    processed_resumes = []
    for uid, files in grouped.items():
        processed_resumes.append({'uid': uid, 'roll': files[0]['roll'],
                                  'name': files[0]['name'], 'latest_date': files[0]['created_at'],
                                  'files': files})
    processed_resumes.sort(key=lambda x: x['roll'])
    db.close()
    return render_template("faculty-department.html", department=dept,
                           students=students, resumes=processed_resumes, session=session)
 
 
@app.route("/debug_dashboard")
def debug_dashboard():
    try:
        db          = get_db()
        departments = db.execute("SELECT DISTINCT department FROM students").fetchall()
        students    = db.execute("SELECT * FROM students LIMIT 3").fetchall()
        return f"""
        <h3>Session</h3><pre>{dict(session)}</pre>
        <h3>Departments</h3><pre>{departments}</pre>
        <h3>Sample Students</h3><pre>{students}</pre>
        """
    except Exception:
        import traceback
        return f"<pre>{traceback.format_exc()}</pre>"
 
 
@app.route("/faculty/edit_student/<uid>", methods=["GET", "POST"])
def faculty_edit_student(uid):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    db = get_db()
    if request.method == "GET":
        row     = db.execute(
            "SELECT s.*, u.email, u.phone FROM students s JOIN users u ON s.uid = u.uid WHERE s.uid = ?",
            [uid]
        ).fetchone()
        student = dict(row) if row else None
        if not student:
            db.close()
            return "Student not found", 404
        db.close()
        return render_template("faculty-edit-student.html", student=student)
 
    department = request.form.get("department")
    roll       = request.form.get("roll")
    name       = request.form.get("name")
    email      = request.form.get("email")
    phone      = request.form.get("phone")
    cgpa       = float(request.form.get("cgpa", 0.0) or 0.0)
    sgpa       = float(request.form.get("sgpa", 0.0) or 0.0)
    attendance = int(request.form.get("attendance", 0) or 0)
 
    if not all([department, roll, name]):
        db.close()
        return "Missing required fields", 400
 
    try:
        db.execute("UPDATE users SET name = ?, email = ?, phone = ? WHERE uid = ?",
                   [name, email, phone, uid])
        db.execute("""
            UPDATE students SET department = ?, roll = ?, name = ?, cgpa = ?, sgpa = ?, attendance = ?
            WHERE uid = ?
        """, [department, roll, name, cgpa, sgpa, attendance, uid])
        db.commit()
    except Exception as e:
        db.rollback()
        return f"Error updating student: {str(e)}", 500
    finally:
        db.close()
    return redirect(url_for("faculty_dashboard"))
 
 
@app.route("/faculty/add_student", methods=["GET"])
def faculty_add_student_form():
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    return render_template("faculty_add_student.html")
 
 
@app.route("/faculty/add_student", methods=["POST"])
def faculty_add_student():
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    username   = request.form.get("username")
    password   = request.form.get("password")
    department = request.form.get("department")
    roll       = request.form.get("roll")
    name       = request.form.get("name")
    email      = request.form.get("email")
    phone      = request.form.get("phone")
    cgpa       = float(request.form.get("cgpa", 0.0) or 0.0)
    sgpa       = float(request.form.get("sgpa", 0.0) or 0.0)
    attendance = int(request.form.get("attendance", 0) or 0)
 
    if not all([username, password, department, roll, name]):
        return "Missing required fields", 400
 
    db = get_db()
    try:
        h = hashlib.sha256(password.encode()).hexdigest()
        db.execute("""
            INSERT INTO users (username, password_hash, name, role, email, phone)
            VALUES (?, ?, ?, 'student', ?, ?)
        """, [username, h, name, email, phone])
        db.commit()
        user_row = db.execute("SELECT uid FROM users WHERE username = ?", [username]).fetchone()
        uid      = user_row[0]
        db.execute("""
            INSERT INTO students (uid, department, roll, name, cgpa, sgpa, attendance)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        """, [uid, department, roll, name, cgpa, sgpa, attendance])
        db.commit()
    except Exception as e:
        db.rollback()
        db.close()
        return str(e)
    return redirect(url_for("faculty_dashboard"))
 
 
@app.route("/faculty/delete_student/<string:roll>", methods=["POST"])
def faculty_delete_student(roll):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    db  = get_db()
    row = db.execute("SELECT uid FROM students WHERE roll = ?", [roll]).fetchone()
    if row:
        uid = row[0]
        db.execute("DELETE FROM resumes WHERE uid = ?", [uid])
        db.execute("DELETE FROM students WHERE roll = ?", [roll])
        db.execute("DELETE FROM users WHERE uid = ?", [uid])
        db.commit()
    db.close()
    return redirect(url_for("faculty_dashboard"))
 
 
# =====================================================================
# FACULTY TO-DO LIST ROUTES
# =====================================================================

@app.route("/add_todo", methods=["POST"])
def add_todo():
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    subject = request.form.get("subject_name")
    department = request.form.get("department_name")
    details = request.form.get("details")
    if subject and department and details:
        db = get_db()
        db.execute(
            "INSERT INTO faculty_todos (faculty_id, subject_name, department_name, details) VALUES (?, ?, ?, ?)",
            [session["uid"], subject, department, details]
        )
        db.commit()
        db.close()
    return redirect(url_for("faculty_dashboard"))

@app.route("/edit_todo/<int:todo_id>", methods=["POST"])
def edit_todo(todo_id):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    subject = request.form.get("subject_name")
    department = request.form.get("department_name")
    details = request.form.get("details")
    if subject and department and details:
        db = get_db()
        db.execute(
            "UPDATE faculty_todos SET subject_name = ?, department_name = ?, details = ? WHERE id = ? AND faculty_id = ?",
            [subject, department, details, todo_id, session["uid"]]
        )
        db.commit()
        db.close()
    return redirect(url_for("faculty_dashboard"))

@app.route("/delete_todo/<int:todo_id>", methods=["POST"])
def delete_todo(todo_id):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    db = get_db()
    db.execute("DELETE FROM faculty_todos WHERE id = ? AND faculty_id = ?", [todo_id, session["uid"]])
    db.commit()
    db.close()
    return redirect(url_for("faculty_dashboard"))

@app.route("/faculty_dashboard")
def faculty_dashboard():
    try:
        uid = session.get("uid")
        if not uid or session.get("role") != "faculty":
            return redirect(url_for("login"))
        dept = request.args.get("dept", None)
        db   = get_db()
 
        def get_students_for(department):
            return db.execute(
                "SELECT s.*, u.email, u.phone FROM students s JOIN users u ON s.uid = u.uid WHERE s.department = ? ORDER BY s.roll",
                [department]
            ).fetchall()
 
        def get_resumes_for(department):
            return db.execute(
                "SELECT r.id, r.uid, r.title, r.created_at, s.roll, s.name FROM resumes r JOIN students s ON r.uid = s.uid WHERE s.department = ? ORDER BY r.created_at DESC",
                [department]
            ).fetchall()
 
        todos = db.execute(
            "SELECT * FROM faculty_todos WHERE faculty_id = ? ORDER BY created_at DESC",
            [uid]
        ).fetchall()

        context = {
            "students_acsml": get_students_for("ACSML"),
            "students_ncsml": get_students_for("NCSML"),
            "students_dcsml": get_students_for("DCSML"),
            "resumes_acsml":  get_resumes_for("ACSML"),
            "resumes_ncsml":  get_resumes_for("NCSML"),
            "resumes_dcsml":  get_resumes_for("DCSML"),
            "todos": todos,
            "session": session,
            "dept": dept,
        }
        db.close()
        return render_template("faculty_dashboard.html", **context)
    except Exception:
        import traceback
        return f"<pre>{traceback.format_exc()}</pre>"
 
 
# =====================================================================
# SIGNUP
# =====================================================================
 
@app.route("/signup", methods=["GET", "POST"])
def signup():
    if request.method == "POST":
        if request.is_json:
            data     = request.get_json()
            name     = data.get("name")
            uid      = data.get("uid")
            email    = data.get("email")
            password = data.get("password")
            role     = data.get("role")
            phone    = data.get("phone", "")
        else:
            name     = request.form.get("name")
            uid      = request.form.get("uid")
            email    = request.form.get("email")
            password = request.form.get("password")
            role     = request.form.get("role")
            phone    = request.form.get("phone", "")
 
        is_hod_val  = 1 if role == "hod" else 0
        db_role_val = "faculty" if role == "hod" else role
        h           = hashlib.sha256(password.encode()).hexdigest()
        db          = get_db()
        try:
            db.execute(
                "INSERT INTO users (uid, username, password_hash, name, role, email, phone, is_hod) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                [uid, uid, h, name, db_role_val, email, phone, is_hod_val]
            )
            if db_role_val == 'student':
                db.execute(
                    "INSERT INTO students (uid, name, department, roll, cgpa, sgpa, attendance) VALUES (?, ?, 'ACSML', ?, 0.0, 0.0, 0)",
                    [uid, name, uid]
                )
            db.commit()
            if request.is_json:
                return jsonify({"success": True, "message": "Account created successfully!"})
            return redirect(url_for('login', message="Account created successfully!"))
        except sqlite3.IntegrityError:
            if request.is_json:
                return jsonify({"success": False, "message": "UID, Username or Email already exists."})
            return render_template("signup.html", error="UID or Email already exists.")
        finally:
            db.close()
    return render_template("signup.html")
 
 
# =====================================================================
# RESUME PDF GENERATION
# =====================================================================
 
def make_styles(template_type):
    styles    = getSampleStyleSheet()
    body_font = "Times-Roman"
    bold_font = "Times-Bold"
    palette   = {
        "classic":  ("#1e3a5f", "#51606f", "#d7dee6"),
        "modern":   ("#0f766e", "#5b6472", "#dbe4ea"),
        "minimal":  ("#111827", "#6b7280", "#e5e7eb"),
        "academic": ("#334155", "#64748b", "#cbd5e1"),
        "creative": ("#8b5e3c", "#647b68", "#e8dfd4"),
    }
    primary, muted, line = palette.get(template_type, palette["classic"])
    primary = colors.HexColor(primary)
    muted   = colors.HexColor(muted)
    line    = colors.HexColor(line)
    return {
        "title":        ParagraphStyle("title", parent=styles["Normal"], fontName=bold_font, fontSize=22, leading=22*1.15, textColor=primary, alignment=TA_LEFT, spaceAfter=3),
        "subtitle":     ParagraphStyle("subtitle", parent=styles["Normal"], fontName=body_font, fontSize=10, leading=10*1.15, textColor=muted, spaceAfter=7),
        "section":      ParagraphStyle("section", parent=styles["Heading2"], fontName=bold_font, fontSize=11.1, leading=11.1*1.15, textColor=primary, spaceBefore=6, spaceAfter=5),
        "body":         ParagraphStyle("body", parent=styles["BodyText"], fontName=body_font, fontSize=10, leading=10*1.15, textColor=colors.black),
        "small":        ParagraphStyle("small", parent=styles["BodyText"], fontName=body_font, fontSize=9, leading=9*1.15, textColor=muted),
        "center_title": ParagraphStyle("center_title", parent=styles["Normal"], fontName=bold_font, fontSize=22, leading=26, textColor=primary, alignment=TA_CENTER, spaceAfter=2),
        "center_sub":   ParagraphStyle("center_sub", parent=styles["Normal"], fontName=body_font, fontSize=10, leading=12, textColor=muted, alignment=TA_CENTER),
        "right_title":  ParagraphStyle("right_title", parent=styles["Normal"], fontName=bold_font, fontSize=18, leading=22, textColor=primary, alignment=TA_RIGHT, spaceAfter=2),
        "justify":      ParagraphStyle("justify", parent=styles["BodyText"], fontName=body_font, fontSize=9.2, leading=12.5, textColor=colors.black, alignment=TA_JUSTIFY),
        "line": line, "primary": primary, "muted": muted,
    }
 
 
def page_header(canvas, doc, form_data, template_type):
    canvas.saveState()
    w, h = A4
    if template_type == "modern":
        canvas.setFillColor(colors.HexColor("#ecfdf5"))
        canvas.rect(0, h - 48, w, 48, stroke=0, fill=1)
        canvas.setStrokeColor(colors.HexColor("#0f766e"))
        canvas.setLineWidth(1.2)
        canvas.line(doc.leftMargin, h - 24, w - doc.rightMargin, h - 24)
        canvas.setFont("Helvetica", 9)
        canvas.setFillColor(colors.HexColor("#0f766e"))
        canvas.drawRightString(w - doc.rightMargin, h - 16, clean_text(form_data.get("field_of_study")))
    elif template_type == "academic":
        canvas.setStrokeColor(colors.HexColor("#94a3b8"))
        canvas.setLineWidth(1)
        canvas.line(doc.leftMargin, h - 22, w - doc.rightMargin, h - 22)
        canvas.setFont("Helvetica", 8.5)
        canvas.setFillColor(colors.HexColor("#64748b"))
        canvas.drawString(doc.leftMargin, h - 16, clean_text(form_data.get("field_of_study")))
    elif template_type == "creative":
        canvas.setFillColor(colors.HexColor("#f7f3ee"))
        canvas.rect(0, h - 52, w, 52, stroke=0, fill=1)
        canvas.setStrokeColor(colors.HexColor("#8b5e3c"))
        canvas.line(doc.leftMargin, h - 26, w - doc.rightMargin, h - 26)
    elif template_type == "minimal":
        canvas.setStrokeColor(colors.HexColor("#111827"))
        canvas.line(25*mm, h - 24, w - 25*mm, h - 24)
    else:
        canvas.setStrokeColor(colors.HexColor("#d7dee6"))
        canvas.line(doc.leftMargin, h - 26, w - doc.rightMargin, h - 26)
        canvas.setFont("Helvetica", 9)
        canvas.setFillColor(colors.HexColor("#51606f"))
        canvas.drawString(doc.leftMargin, h - 18, clean_text(form_data.get("field_of_study")))
    canvas.restoreState()
 
 
def box(title, text, styles, fill=None):
    tbl = Table([[Paragraph(f"<b>{title}</b>", styles["body"])],
                 [Paragraph(text, styles["small"])]], colWidths=[None])
    ts  = [
        ("BOX",           (0, 0), (-1, -1), 0.8, styles["line"]),
        ("VALIGN",        (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 6),
        ("TOPPADDING",    (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]
    if fill:
        ts.append(("BACKGROUND", (0, 0), (-1, -1), fill))
    tbl.setStyle(TableStyle(ts))
    return tbl
 
 
def common_sections(form_data, styles):
    story = []
    story.append(Paragraph("CAREER OBJECTIVE", styles["section"]))
    story.append(Paragraph(enhance_objective(form_data.get("objective"), form_data.get("field_of_study")), styles["body"]))
    story.append(Spacer(1, 5))
    story.append(Paragraph("ACADEMIC QUALIFICATION", styles["section"]))
    if form_data.get("ssc_school"):
        story.append(Paragraph(f"<b>SSC (10th Standard)</b>: {clean_text(form_data.get('ssc_school'))}<br/>Score: {clean_text(form_data.get('ssc_percentage'))}", styles["body"]))
        story.append(Spacer(1, 3))
    if form_data.get("inter_college"):
        story.append(Paragraph(f"<b>Intermediate (12th Standard)</b>: {clean_text(form_data.get('inter_college'))}<br/>Score: {clean_text(form_data.get('inter_percentage'))}", styles["body"]))
        story.append(Spacer(1, 3))
    if form_data.get("ug_college"):
        story.append(Paragraph(f"<b>UG</b>: {clean_text(form_data.get('ug_degree'))} in {clean_text(form_data.get('ug_branch'))}<br/>{clean_text(form_data.get('ug_college'))}<br/>Score: {clean_text(form_data.get('ug_score'))}", styles["body"]))
        story.append(Spacer(1, 3))
    if form_data.get("pg_college"):
        story.append(Paragraph(f"<b>PG</b>: {clean_text(form_data.get('pg_degree'))} in {clean_text(form_data.get('pg_branch'))}<br/>{clean_text(form_data.get('pg_college'))}<br/>Score: {clean_text(form_data.get('pg_score'))}", styles["body"]))
        story.append(Spacer(1, 3))
    skill_names  = form_data.get("skill_name", [])
    skill_levels = form_data.get("skill_level", [])
    if isinstance(skill_names, str):
        skill_names  = [skill_names]
        skill_levels = [skill_levels]
    processed_skills = []
    for i in range(len(skill_names)):
        name  = clean_text(skill_names[i])
        level = clean_text(skill_levels[i]) if i < len(skill_levels) else ""
        if name:
            processed_skills.append(f"{name} - {level}")
    form_data['processed_technical_skills'] = " • ".join(processed_skills)
    return story

 
def parse_projects_with_bullets(form_data):

    titles = form_data.get("project_title", [])
    descriptions = form_data.get("project_description", [])
    technologies = form_data.get("project_technology", [])

    if not isinstance(titles, list):
        titles = [titles]

    if not isinstance(descriptions, list):
        descriptions = [descriptions]

    if not isinstance(technologies, list):
        technologies = [technologies]

    projects = []

    for title, tech, desc in zip(
        titles,
        technologies,
        descriptions
    ):

        bullets = [
            line.strip()
            for line in str(desc).split("\n")
            if line.strip()
        ]

        projects.append({
            "title": str(title),
            "technology": str(tech),
            "bullets": bullets
        })

    return projects
 
 
def generate_resume_pdf(template_type, form_data):
    buffer        = io.BytesIO()
    template_type = (template_type or "classic").lower()
    if template_type not in ["classic", "modern", "minimal", "academic", "creative"]:
        template_type = "classic"
    styles = make_styles(template_type)
 
    def on_page(canvas, doc):
        page_header(canvas, doc, form_data, template_type)
        canvas.saveState()
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(colors.HexColor("#6b7280"))
        canvas.drawRightString(A4[0] - doc.rightMargin, 8*mm, f"Page {doc.page}")
        canvas.restoreState()
 
    if template_type == "modern":
        doc      = BaseDocTemplate(buffer, pagesize=A4, leftMargin=12*mm, rightMargin=12*mm, topMargin=54*mm, bottomMargin=15*mm)
        sidebar_w = 60*mm
        gap       = 6*mm
        main_w    = A4[0] - doc.leftMargin - doc.rightMargin - sidebar_w - gap
        frame_h   = A4[1] - doc.topMargin - doc.bottomMargin
        left_frame  = Frame(doc.leftMargin, doc.bottomMargin, sidebar_w, frame_h, id="sidebar", showBoundary=0)
        right_frame = Frame(doc.leftMargin + sidebar_w + gap, doc.bottomMargin, main_w, frame_h, id="main", showBoundary=0)
        doc.addPageTemplates([PageTemplate(id="modern", frames=[left_frame, right_frame], onPage=on_page)])
        story = []
        story.append(Paragraph("PROFILE", styles["section"]))
        story.append(Paragraph(enhance_objective(form_data.get("objective"), form_data.get("field_of_study")), styles["small"]))
        story.append(Spacer(1, 4))
        story.append(Paragraph("CONTACT", styles["section"]))
        contact_lines = [f"Email: {clean_text(form_data.get('email'))}", f"Phone: {clean_text(form_data.get('phone'))}"]
        if form_data.get('linkedin'): contact_lines.append(f"LinkedIn: {clean_text(form_data.get('linkedin'))}")
        if form_data.get('github'):   contact_lines.append(f"GitHub: {clean_text(form_data.get('github'))}")
        if form_data.get('website'):  contact_lines.append(f"Website: {clean_text(form_data.get('website'))}")
        location_str = ", ".join(filter(None, [clean_text(form_data.get('city')), clean_text(form_data.get('state')), clean_text(form_data.get('country'))]))
        if location_str: contact_lines.append(f"Location: {location_str}")
        story.append(Paragraph("<br/>".join(contact_lines), styles["small"]))
        story.append(Spacer(1, 4))
        story.append(Paragraph("SKILLS", styles["section"]))
        skill_names  = form_data.get("skill_name", [])
        skill_levels = form_data.get("skill_level", [])
        if isinstance(skill_names, str): skill_names = [skill_names]
        if isinstance(skill_levels, str): skill_levels = [skill_levels]
        processed_skills = []
        for i in range(len(skill_names)):
            name  = clean_text(skill_names[i])
            level = skill_levels[i] if i < len(skill_levels) else ""
            if name: processed_skills.append(f"{name} - {level}")
        story.append(Paragraph(" • ".join(processed_skills) or "N/A", styles["small"]))
        story.append(Spacer(1, 4))
        story.append(Paragraph("LANGUAGES", styles["section"]))
        langs = "<br/>".join(f"• {clean_text(x)}" for x in split_items(form_data.get("languages")))
        story.append(Paragraph(langs or "N/A", styles["small"]))
        story.append(FrameBreak())
        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["subtitle"]))
        story.extend(common_sections(form_data, styles))
        projects = parse_projects_with_bullets(form_data)
        if projects:
            story.append(Paragraph("PROJECTS", styles["section"]))
            for i, p in enumerate(projects, 1):
                title   = clean_text(p["title"]) or f"Project {i}"
                bullets = p["bullets"] or ["Project details not provided."]
                txt     = f"<b>{title}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets)
                story.append(box(f"Project {i}", txt, styles, fill=colors.HexColor("#f0fdfa")))
                story.append(Spacer(1, 4))
        experiences = split_items(form_data.get("work_experience")) or split_items(form_data.get("internships"))
        if experiences:
            story.append(Paragraph("INTERNSHIPS", styles["section"]))
            for exp in experiences:
                bullets_exp = describe_experience(exp)
                txt         = f"<b>{clean_text(exp)}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets_exp)
                story.append(Paragraph(txt, styles["body"]))
                story.append(Spacer(1, 4))
        certs = split_items(form_data.get("certifications"))
        if certs:
            story.append(Paragraph("CERTIFICATIONS", styles["section"]))
            for c in certs:
                story.append(Paragraph(f"• <b>{clean_text(c)}</b>: {clean_text(describe_certification(c))}", styles["body"]))
                story.append(Spacer(1, 2))
        doc.build(story)
        pdf_out = buffer.getvalue()
        buffer.close()
        return pdf_out
 
    # All other templates
    doc   = BaseDocTemplate(buffer, pagesize=A4, leftMargin=15*mm, rightMargin=15*mm, topMargin=18*mm, bottomMargin=15*mm)
    story = []
    if template_type == "classic":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        contact_parts = [clean_text(form_data.get('field_of_study')), clean_text(form_data.get('email')), clean_text(form_data.get('phone'))]
        if form_data.get('linkedin'): contact_parts.append(f"LinkedIn: {clean_text(form_data.get('linkedin'))}")
        if form_data.get('github'):   contact_parts.append(f"GitHub: {clean_text(form_data.get('github'))}")
        if form_data.get('website'):  contact_parts.append(f"Website: {clean_text(form_data.get('website'))}")
        story.append(Paragraph(" | ".join(filter(None, contact_parts)), styles["subtitle"]))
        story.append(HRFlowable(width="100%", thickness=0.8, color=styles["line"]))
        story.append(Spacer(1, 5))
    elif template_type == "minimal":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["center_title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["center_sub"]))
        contact_parts = [clean_text(form_data.get('email')), clean_text(form_data.get('phone'))]
        if form_data.get('linkedin'): contact_parts.append(clean_text(form_data.get('linkedin')))
        if form_data.get('github'):   contact_parts.append(clean_text(form_data.get('github')))
        story.append(Paragraph(" • ".join(filter(None, contact_parts)), styles["center_sub"]))
        story.append(Spacer(1, 5))
        story.append(HRFlowable(width="42%", thickness=1, color=styles["line"]))
        story.append(Spacer(1, 7))
    elif template_type == "academic":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["right_title"]))
        contact_info = f"{clean_text(form_data.get('email'))} | {clean_text(form_data.get('phone'))}"
        if form_data.get('linkedin'): contact_info += f" | LinkedIn: {clean_text(form_data.get('linkedin'))}"
        story.append(Paragraph(contact_info, styles["small"]))
        story.append(Spacer(1, 4))
        story.append(HRFlowable(width="100%", thickness=1, color=styles["line"]))
        story.append(Spacer(1, 6))
    elif template_type == "creative":
        story.append(Paragraph(clean_text(form_data.get("name")), styles["title"]))
        story.append(Paragraph(clean_text(form_data.get("field_of_study")), styles["subtitle"]))
        if form_data.get('website'): story.append(Paragraph(f"Portfolio: {clean_text(form_data.get('website'))}", styles["small"]))
        story.append(Spacer(1, 3))
        story.append(HRFlowable(width="100%", thickness=1.1, color=styles["line"]))
        story.append(Spacer(1, 7))
 
    story.extend(common_sections(form_data, styles))
 
    skill_names  = form_data.get("skill_name", [])
    skill_levels = form_data.get("skill_level", [])
    if isinstance(skill_names, str): skill_names = [skill_names]
    if isinstance(skill_levels, str): skill_levels = [skill_levels]
    processed_skills = []
    for i in range(len(skill_names)):
        name  = clean_text(skill_names[i])
        level = skill_levels[i] if i < len(skill_levels) else ""
        if name: processed_skills.append(f"{name} - {level}")
    if processed_skills:
        story.append(Paragraph("TECHNICAL SKILLS", styles["section"]))
        story.append(Paragraph(" • ".join(processed_skills), styles["body"]))
        story.append(Spacer(1, 5))
 
    if form_data.get("languages"):
        story.append(Paragraph("LANGUAGES", styles["section"]))
        if template_type == "minimal":
            data = [[Paragraph("Languages", styles["body"]), Paragraph(" • ".join(clean_text(x) for x in split_items(form_data.get("languages"))), styles["small"])]]
            t    = Table(data, colWidths=[35*mm, None])
            t.setStyle(TableStyle([("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#f8fafc")), ("BOX", (0, 0), (-1, -1), 0.8, styles["line"]), ("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6), ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5)]))
            story.append(t)
        elif template_type == "academic":
            story.append(Paragraph(" • ".join(clean_text(x) for x in split_items(form_data.get("languages"))), styles["justify"]))
        else:
            story.append(Paragraph(" • ".join(clean_text(x) for x in split_items(form_data.get("languages"))), styles["body"]))
        story.append(Spacer(1, 5))
 
    projects = parse_projects_with_bullets(form_data)
    if projects:
        story.append(Paragraph("PROJECTS", styles["section"]))
        for i, p in enumerate(projects, 1):
            title   = clean_text(p["title"]) or f"Project {i}"
            bullets = p["bullets"] or ["Project details not provided."]
            txt     = f"<b>{title}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets)
            if template_type == "creative":
                story.append(box(f"Project {i}", txt, styles, fill=colors.HexColor("#fbf7f2")))
            elif template_type == "academic":
                story.append(Paragraph(txt, styles["justify"]))
            else:
                story.append(Paragraph(txt, styles["body"]))
            story.append(Spacer(1, 4))
 
    experiences = split_items(form_data.get("work_experience")) or split_items(form_data.get("internships"))
    if experiences:
        story.append(Paragraph("EXPERIENCE & INTERNSHIPS", styles["section"]))
        for exp in experiences:
            bullets_exp = describe_experience(exp)
            txt         = f"<b>{clean_text(exp)}</b><br/>" + "<br/>".join(f"• {clean_text(b)}" for b in bullets_exp)
            story.append(Paragraph(txt, styles["body"]))
            story.append(Spacer(1, 5))
        story.append(Spacer(1, 5))
 
    certs = split_items(form_data.get("certifications"))
    if certs:
        story.append(Paragraph("CERTIFICATIONS", styles["section"]))
        for c in certs:
            story.append(Paragraph(f"• <b>{clean_text(c)}</b>: {clean_text(describe_certification(c))}", styles["body"]))
            story.append(Spacer(1, 3))
 
    doc.addPageTemplates([PageTemplate(id=template_type, frames=[Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="main")], onPage=on_page)])
    doc.build(story)
    pdf_out = buffer.getvalue()
    buffer.close()
    return pdf_out
 
 
# =====================================================================
# DOWNLOAD / VIEW RESUME
# =====================================================================
 
@app.route("/download_resume/<int:resume_id>")
def download_resume(resume_id):
    db  = get_db()
    row = db.execute("SELECT pdf_content, title FROM resumes WHERE id = ?", [resume_id]).fetchone()
    db.close()
    if not row:
        return "Resume not found in database", 404
    filename = f"{row['title'].replace(' ', '_')}.pdf"
    return send_file(io.BytesIO(row["pdf_content"]), as_attachment=True, download_name=filename, mimetype="application/pdf")
 
 
@app.route("/view_resume/<int:resume_id>")
def view_resume(resume_id):
    db = get_db()
    if session.get("role") == "faculty":
        row = db.execute("SELECT pdf_content FROM resumes WHERE id = ?", [resume_id]).fetchone()
    else:
        row = db.execute("SELECT pdf_content FROM resumes WHERE id = ? AND uid = ?",
                         [resume_id, session.get("uid")]).fetchone()
    db.close()
    if not row:
        return "Resume not found", 404
    return send_file(io.BytesIO(row["pdf_content"]), mimetype="application/pdf")
 
 
@app.route("/faculty/resumes")
def faculty_resumes():
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    db   = get_db()
    rows = db.execute("""
        SELECT r.id, r.uid, r.title, r.created_at, s.roll, s.name, s.department
        FROM resumes r
        JOIN students s ON r.uid = s.uid
        ORDER BY s.roll, r.created_at DESC
    """).fetchall()
    from collections import defaultdict
    grouped = defaultdict(list)
    for r in rows:
        grouped[r['uid']].append(dict(r))
    resumes = []
    for uid, st_resumes in grouped.items():
        resumes.append({'uid': uid, 'roll': st_resumes[0]['roll'], 'name': st_resumes[0]['name'],
                        'department': st_resumes[0]['department'], 'latest_date': st_resumes[0]['created_at'],
                        'files': st_resumes})
    resumes.sort(key=lambda x: x['latest_date'], reverse=True)
    db.close()
    return render_template("faculty_resumes.html", resumes=resumes, session=session)
 
 
@app.route("/combined_resume/<uid>")
def combined_resume(uid):
    if not session.get("uid"):
        return redirect(url_for("login"))
    if session.get("role") != "faculty" and session.get("uid") != uid:
        return "Unauthorized", 403
    db   = get_db()
    rows = db.execute("SELECT pdf_content FROM resumes WHERE uid = ? ORDER BY created_at ASC", [uid]).fetchall()
    db.close()
    if not rows:
        return "No resumes found", 404
    if len(rows) == 1:
        return send_file(io.BytesIO(rows[0]["pdf_content"]), mimetype="application/pdf")
    try:
        from pypdf import PdfWriter
        merger = PdfWriter()
        for row in rows:
            if row["pdf_content"]:
                merger.append(io.BytesIO(row["pdf_content"]))
        output = io.BytesIO()
        merger.write(output)
        output.seek(0)
        return send_file(output, as_attachment=True, download_name=f"Combined_Resumes_{uid}.pdf", mimetype="application/pdf")
    except ImportError:
        return send_file(io.BytesIO(rows[-1]["pdf_content"]), mimetype="application/pdf")
 
 
@app.route("/delete_resume", methods=["POST"])
def delete_resume():
    if "uid" not in session:
        return redirect(url_for("login"))
    resume_id_str = request.form.get("resume_id")
    if not resume_id_str:
        return "Invalid request", 400
    try:
        resume_id = int(resume_id_str)
    except ValueError:
        return "Invalid resume ID", 400
    uid = session["uid"]
    db  = get_db()
    db.execute("DELETE FROM resumes WHERE id = ? AND uid = ?", [resume_id, uid])
    db.commit()
    db.close()
    return redirect(url_for("student_dashboard"))
 
 
# =====================================================================
# STUDENT DETAILS
# =====================================================================
 
@app.route("/student_details/<student_uid>")
def student_details(student_uid):
    if session.get("role") != "faculty":
        return redirect(url_for("login"))
    db      = get_db()
    student = db.execute("""
        SELECT s.uid, s.roll, s.name, s.department, s.cgpa, s.sgpa, s.attendance, u.email, u.phone
        FROM students s
        LEFT JOIN users u ON s.uid = u.uid
        WHERE s.uid = ?
    """, [student_uid]).fetchone()
    if not student:
        db.close()
        return "Student not found", 404
    resumes = db.execute("SELECT * FROM resumes WHERE uid = ? ORDER BY created_at DESC", [student_uid]).fetchall()
    db.close()
    return render_template("student_details.html", student=student, resumes=resumes)
 
 
# =====================================================================
# GEMINI PROJECT DESCRIPTION
# =====================================================================
 
@app.route("/generate_project_description", methods=["POST"])
def generate_project_description():
    data     = request.get_json()
    title    = data.get("title", "")
    features = data.get("features", "")
    prompt   = f"""
Project Title: {title}
 
Features:
{features}
 
Generate ONLY 5 resume bullet points.
 
Rules:
- Start each bullet with •
- ATS friendly
- Maximum 20 words per bullet
- Use action verbs
- No headings, No numbering, No explanations, No markdown
- Return only the bullet points
"""
    try:
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        return jsonify({"success": True, "description": response.text})
    except Exception as e:
        logger.error(f"Gemini Error: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500
 
 
# =====================================================================
# NOTES
# =====================================================================
 
@app.route('/upload_notes', methods=['GET', 'POST'])
def upload_notes():
    uid = session.get('uid')
    if not uid:
        return redirect('/login')
    if request.method == 'POST':
        title         = request.form.get('title')
        description   = request.form.get('description')
        department_id = request.form.get('department_id')
        file          = request.files.get('file')
        if not file or file.filename == '':
            flash('Please select a file', 'error')
            return redirect(url_for('upload_notes'))
        filename        = secure_filename(file.filename)
        timestamp       = datetime.now().strftime('%Y%m%d_%H%M%S')
        unique_filename = f"{timestamp}_{filename}"
        filepath        = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        file.save(filepath)
        conn = get_db()
        conn.execute("""
            INSERT INTO notes (title, description, file_path, file_name, department_id, faculty_id)
            VALUES (?, ?, ?, ?, ?, ?)
        """, (title, description, unique_filename, filename, department_id, uid))
        conn.commit()
        conn.close()
        flash('Notes uploaded successfully!', 'success')
        return redirect(url_for('faculty_my_notes'))
    conn        = get_db()
    departments = conn.execute("SELECT * FROM departments ORDER BY name").fetchall()
    conn.close()
    return render_template('upload_notes.html', departments=departments)
 
 
@app.route('/faculty_my_notes')
def faculty_my_notes():
    uid = session.get('uid')
    if not uid:
        return redirect('/login')
    conn  = get_db()
    notes = conn.execute('''
        SELECT n.id, n.title, n.description, n.file_name, n.uploaded_at, d.name AS department_name
        FROM notes n
        JOIN departments d ON n.department_id = d.id
        WHERE n.faculty_id = ?
        ORDER BY n.uploaded_at DESC
    ''', (uid,)).fetchall()
    conn.close()
    return render_template('faculty_my_notes.html', notes=notes)
 
 
# ── FIX 11: Added /view_notes route (referenced in base.html sidebar) ─────────
@app.route('/view_notes')
def view_notes():
    uid = session.get('uid')
    if not uid or session.get('role') != 'student':
        return redirect(url_for('login'))
    db  = get_db()
    # Find student's department
    student = db.execute("SELECT department FROM students WHERE uid = ?", [uid]).fetchone()
    dept_name = student['department'] if student else None
 
    department = None
    notes      = []
    if dept_name:
        department = db.execute("SELECT * FROM departments WHERE name = ?", [dept_name]).fetchone()
        if department:
            notes = db.execute('''
                SELECT n.id, n.title, n.description, n.file_name, n.file_path, n.uploaded_at,
                       u.name AS faculty_name
                FROM notes n
                JOIN users u ON n.faculty_id = u.uid
                WHERE n.department_id = ?
                ORDER BY n.uploaded_at DESC
            ''', (department['id'],)).fetchall()
    db.close()
    return render_template('view_notes.html', notes=notes, department=department)
 
 
# ── FIX 12: Added /download_note route (used in view_notes.html) ─────────────
@app.route('/download_note/<filename>')
def download_note(filename):
    note_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.exists(note_path):
        return "File not found", 404
    return send_file(note_path, as_attachment=True)
 
 
@app.route('/delete_note/<int:note_id>')
def delete_note(note_id):
    uid = session.get('uid')
    if not uid:
        return redirect('/login')
    conn = get_db()
    note = conn.execute("SELECT * FROM notes WHERE id = ? AND faculty_id = ?", (note_id, uid)).fetchone()
    if note:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], note['file_path'])
        if os.path.exists(file_path):
            os.remove(file_path)
        conn.execute("DELETE FROM notes WHERE id = ?", (note_id,))
        conn.commit()
    conn.close()
    flash('Note deleted successfully!', 'success')
    return redirect(url_for('faculty_my_notes'))
 
 
# =====================================================================
# MISC
# =====================================================================
 
@app.route("/check-auth")
def check_auth_endpoint():
    return jsonify({"logged_in": "uid" in session, "role": session.get("role")})
 
 
@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))
 
 
@app.context_processor
def inject_recent_resumes():
    if "uid" in session and session.get("role") == "student":
        try:
            db     = get_db()
            rows   = db.execute(
                "SELECT id, title FROM resumes WHERE uid = ? ORDER BY created_at DESC LIMIT 3",
                [session["uid"]]
            ).fetchall()
            recent = [dict(r) for r in rows]
            db.close()
            return dict(recent_resumes=recent)
        except Exception:
            return dict(recent_resumes=[])
    return dict(recent_resumes=[])

#----------------Admin---------------------------

@app.route("/manage_students")
def manage_students():

    if session.get("role") != "admin":
        return redirect(url_for("login"))

    return render_template("manage_students.html")

@app.route("/manage_faculty")
def manage_faculty():

    if session.get("role") != "admin":
        return redirect(url_for("login"))

    conn = get_db()

    faculty = conn.execute(
        "SELECT * FROM users WHERE role='faculty'"
    ).fetchall()

    conn.close()

    return render_template(
        "manage_faculty.html",
        faculty=faculty
    )

@app.route("/faculty_permissions")
def faculty_permissions():

    if session.get("role") != "admin":
        return redirect(url_for("login"))

    conn = get_db()

    faculty = conn.execute("""
        SELECT *
        FROM users
        WHERE role='faculty'
    """).fetchall()

    conn.close()

    return render_template(
        "faculty_permissions.html",
        faculty=faculty
    )




@app.route("/test-ai")
def test_ai():

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents="Say Hello"
    )

    return response.text
# =====================================================================
# STARTUP
# =====================================================================
 
if __name__ == "__main__":

    init_database()

    port = int(os.environ.get("PORT", 5000))

    app.run(
        debug=os.environ.get("FLASK_DEBUG", "true").lower() == "true",
        host="0.0.0.0",
        port=port
    )